"""사용자 PPT에 일괄 등록 슬라이드 2장 더 추가 (총 11 → 13장).
- 슬라이드 12: 📥 일괄 등록 — 입구 (샘플 엑셀 + 드롭존)
- 슬라이드 13: ✅ 업로드 후 — 5문항 표 + 미리보기
"""
import os, sys, shutil
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = lambda f: os.path.join('data/presentation/img', f)
PPT = 'data/presentation/엑셀-등록페이지_연동.pptx'

FONT = '맑은 고딕'
MONO = 'Consolas'

C_PRIMARY = RGBColor(0xC2, 0x18, 0x5B)
C_DARK = RGBColor(0x1F, 0x29, 0x37)
C_MUTED = RGBColor(0x6B, 0x72, 0x80)
C_BG = RGBColor(0xFA, 0xFA, 0xFA)
C_OK = RGBColor(0x05, 0x96, 0x69)
C_WARN = RGBColor(0xD9, 0x77, 0x06)
C_VIDEO = RGBColor(0x7C, 0x3A, 0xED)
C_PINK = RGBColor(0xDB, 0x27, 0x77)
C_META = RGBColor(0x47, 0x55, 0x69)
C_ACCENT = RGBColor(0x21, 0x96, 0xF3)


def tb(slide, l, t, w, h, lines, *, size=14, bold=False, color=C_DARK,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return box


def rect(slide, l, t, w, h, fill, line=None, line_w=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def title_bar(slide, prs, t, sub=None):
    rect(slide, 0, 0, prs.slide_width, Inches(0.15), C_PRIMARY)
    tb(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
       t, size=26, bold=True)
    if sub:
        tb(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4),
           sub, size=13, color=C_MUTED)


prs = Presentation(PPT)
BLANK = prs.slide_layouts[6]
print(f'기존 슬라이드 {len(prs.slides)}장 — 신규 2장 append')


# =============================================================
# Slide 12 — 📥 일괄 등록 — 입구 (샘플 엑셀 + 드롭존)
# =============================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, prs, '9. 📥 일괄 등록 — 엑셀 한 파일에 여러 문항 + 강의주소',
          '운영자가 엑셀에 N개 문항·강의주소·블록 시퀀스 모두 채워서 한 번에 업로드')

# 좌: 샘플 엑셀 캡처
tb(s, Inches(0.4), Inches(1.4), Inches(7.2), Inches(0.4),
   '① 샘플 엑셀 — 5문항 + 강의주소 (sample_batch_5문항.xlsx)',
   size=14, bold=True, color=C_META)
s.shapes.add_picture(IMG('xlsx_sample_batch_5문항_p1.png'),
                     Inches(0.4), Inches(1.85), width=Inches(7.2))
tb(s, Inches(0.4), Inches(5.7), Inches(7.2), Inches(0.4),
   '· 보라색 「강의」 영역에 YouTube URL 5건 모두 채워짐  · 1·21·51번 + 추가 2문항',
   size=10, color=C_MUTED, align=PP_ALIGN.CENTER)

# 우: 일괄 등록 페이지 빈 화면 (드롭존)
tb(s, Inches(7.8), Inches(1.4), Inches(5.2), Inches(0.4),
   '② 일괄 등록 페이지 — 드롭존',
   size=14, bold=True, color=C_PINK)
s.shapes.add_picture(IMG('batch_empty.png'),
                     Inches(7.8), Inches(1.85), width=Inches(5.2))

# 하단 흐름 설명
rect(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(1.0),
     RGBColor(0xFC, 0xE7, 0xF3), line=C_PINK, line_w=2)
tb(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
   '🔄 입구 흐름', size=13, bold=True, color=C_PINK)
tb(s, Inches(0.6), Inches(6.75), Inches(12.2), Inches(0.6),
   [
       '1) 운영자가 v3 양식 엑셀에 N개 문항 작성 (강의주소·해설 블록 시퀀스 포함)',
       '2) 일괄 등록 페이지(/admin/register-batch)에 파일 드롭 또는 클릭 선택',
       '3) /api/bulk-parse가 시트 파싱 → 블록 시퀀스 분리 → JSON 배열 반환',
   ], size=11)


# =============================================================
# Slide 13 — ✅ 업로드 후 — 표 + 미리보기
# =============================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, prs, '10. ✅ 업로드 후 — 5문항 검수 + 한 번에 등록',
          '좌측: 문항 리스트 표 / 우측: 행 클릭 시 실시간 미리보기')

# 전체 캡처 (페이지 풀)
s.shapes.add_picture(IMG('batch_uploaded.png'),
                     Inches(0.3), Inches(1.4), width=Inches(8.5))

# 우측: 동작 설명 카드
right = Inches(9.0)
rect(s, right, Inches(1.4), Inches(4.0), Inches(5.8), C_BG, line=C_OK, line_w=2)
tb(s, right + Inches(0.15), Inches(1.55), Inches(3.7), Inches(0.4),
   '⚙ 동작', size=14, bold=True, color=C_OK)

tb(s, right + Inches(0.15), Inches(2.0), Inches(3.7), Inches(5.0),
   [
       '◆ 상단 요약',
       '   파일명 · 총 N문항',
       '   🎬 강의주소 N건',
       '   ⚠ 그림 마커 N건',
       '',
       '◆ 표 컬럼',
       '   # · 문항코드 · 과목 ·',
       '   발문 요약 · 정답 ·',
       '   🎬 (강의 유무) · 상태',
       '',
       '◆ 상태 컬럼',
       '   ✓ 정상',
       '   ⚠ 그림 마커 미해결',
       '',
       '◆ 행 클릭',
       '   우측에 미리보기 갱신',
       '   (강의 버튼·블록 시퀀스',
       '    모두 그대로 렌더링)',
       '',
       '◆ 액션',
       '   [JSON 다운로드]',
       '   [전체 등록] — 백엔드 미정',
   ], size=11)

# 하단: 가치 한 줄
rect(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
     RGBColor(0xD1, 0xFA, 0xE5))
tb(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
   '✅ 엑셀 한 번 업로드 = 5문항 + 강의주소 + 블록 시퀀스 + 그림 마커 모두 일괄 등록 (검수 후)',
   size=12, color=C_OK, bold=True, align=PP_ALIGN.CENTER)


try:
    prs.save(PPT)
    print(f'WROTE: {PPT}')
except PermissionError:
    alt = PPT.replace('.pptx', '_v2.pptx')
    prs.save(alt)
    print(f'(원본 잠김) WROTE: {alt}')
print(f'  슬라이드 {len(prs.slides)}장')
