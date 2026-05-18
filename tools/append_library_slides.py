"""PPT에 시험 페이지 슬라이드 2장 더 추가 (총 13 → 15장).
- 슬라이드 14: 📚 시험 문항 라이브러리 — 메타 4·5와 연결
- 슬라이드 15: 🔍 필터 / 정렬 동작 예시
"""
import os, sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = lambda f: os.path.join('data/presentation/img', f)
PPT_INPUT = 'data/presentation/엑셀-등록페이지_연동_v2.pptx'
PPT_OUTPUT = 'data/presentation/엑셀-등록페이지_연동_v3.pptx'

FONT = '맑은 고딕'
MONO = 'Consolas'

C_PRIMARY = RGBColor(0xC2, 0x18, 0x5B)
C_DARK = RGBColor(0x1F, 0x29, 0x37)
C_MUTED = RGBColor(0x6B, 0x72, 0x80)
C_BG = RGBColor(0xFA, 0xFA, 0xFA)
C_OK = RGBColor(0x05, 0x96, 0x69)
C_VIDEO = RGBColor(0x7C, 0x3A, 0xED)
C_CYAN = RGBColor(0x08, 0x91, 0xB2)
C_ROSE = RGBColor(0xE1, 0x1D, 0x48)
C_INDIGO = RGBColor(0x4F, 0x46, 0xE5)
C_ACCENT = RGBColor(0x21, 0x96, 0xF3)


def tb(slide, l, t, w, h, lines, *, size=14, bold=False, color=C_DARK,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return box


def rect(slide, l, t, w, h, fill, line=None, line_w=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def title_bar(slide, prs, t, sub=None):
    rect(slide, 0, 0, prs.slide_width, Inches(0.15), C_PRIMARY)
    tb(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
       t, size=26, bold=True)
    if sub:
        tb(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4),
           sub, size=13, color=C_MUTED)


prs = Presentation(PPT_INPUT)
BLANK = prs.slide_layouts[6]
print(f'기존 슬라이드 {len(prs.slides)}장 — 신규 2장 append')


# =============================================================
# Slide 14 — 📚 시험 문항 라이브러리 (메타 슬라이드 4·5와 연결)
# =============================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, prs, '11. 📚 시험 페이지 — 문항 라이브러리 (등록된 문항 배열)',
          '슬라이드 4·5에서 정의한 분류(세로축) × 속성(가로축)을 실제 UI로 — 정렬 + 필터')

# 풀 캡처
s.shapes.add_picture(IMG('library_default.png'),
                     Inches(0.3), Inches(1.4), width=Inches(9.0))

# 우측 설명
right = Inches(9.6)
rect(s, right, Inches(1.4), Inches(3.5), Inches(5.5), C_BG, line=C_CYAN, line_w=2)
tb(s, right + Inches(0.15), Inches(1.55), Inches(3.2), Inches(0.4),
   '🗂 좌측 — 분류 트리 (세로축)', size=13, bold=True, color=C_CYAN)
tb(s, right + Inches(0.15), Inches(1.95), Inches(3.2), Inches(1.6),
   [
       '· 과목ID > 챕터 트리',
       '· 클릭하면 해당 카테고리만 필터',
       '· 문항 수 카운트 표시',
       '· 전체 보기 ↔ 카테고리 전환',
   ], size=10)

tb(s, right + Inches(0.15), Inches(3.55), Inches(3.2), Inches(0.4),
   '🏷 상단 — 속성 필터 (가로축)', size=13, bold=True, color=C_ROSE)
tb(s, right + Inches(0.15), Inches(3.95), Inches(3.2), Inches(1.6),
   [
       '· 연도 / 회차 / 문제유형',
       '· 빈출도 / 난이도 ★ 이상',
       '· 발문·분류 검색',
       '· [초기화] 한 번에 리셋',
   ], size=10)

tb(s, right + Inches(0.15), Inches(5.55), Inches(3.2), Inches(0.4),
   '📋 메인 — 표 + 정렬', size=13, bold=True, color=C_INDIGO)
tb(s, right + Inches(0.15), Inches(5.95), Inches(3.2), Inches(0.9),
   [
       '· 컬럼 헤더 클릭 = 정렬',
       '· 문항코드 / 과목 / 출처 / 빈출 / 난이도',
   ], size=10)

# 하단: 데이터 흐름
rect(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
     RGBColor(0xCF, 0xFA, 0xFE))
tb(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
   '🔄 데이터 소스: 일괄 등록한 sample_batch_5문항.xlsx → src/data/questions.json → 페이지 정적 import',
   size=11, color=C_CYAN, bold=True, align=PP_ALIGN.CENTER)


# =============================================================
# Slide 15 — 🔍 필터 / 정렬 동작 예시
# =============================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, prs, '12. 🔍 시험 페이지 — 필터·정렬 동작 예시',
          '같은 5문항을 두 가지 다른 조건으로 본 화면 비교')

# 좌측: 과목 = 전기기기만
tb(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(0.4),
   '① 분류 필터 — 「전기기기」만', size=14, bold=True, color=C_CYAN)
s.shapes.add_picture(IMG('library_subject_machinery.png'),
                     Inches(0.4), Inches(1.85), width=Inches(6.0))

# 우측: 21번 차동계전기 선택 — 강의 버튼 + 블록 시퀀스 미리보기
tb(s, Inches(6.8), Inches(1.4), Inches(6.2), Inches(0.4),
   '② 21번 차동계전기 선택 — 강의 버튼 + 해설 블록 미리보기',
   size=14, bold=True, color=C_VIDEO)
s.shapes.add_picture(IMG('library_select_q21.png'),
                     Inches(6.8), Inches(1.85), width=Inches(6.0))

# 하단: 가치
rect(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
     RGBColor(0xD1, 0xFA, 0xE5))
tb(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
   '✅ 운영자 등록 → 학습자가 분류·속성·강의·정답·해설을 그대로 볼 수 있는 카탈로그까지 한 흐름',
   size=12, color=C_OK, bold=True, align=PP_ALIGN.CENTER)


try:
    prs.save(PPT_INPUT)   # 같은 v2 파일에 덮어쓰기
    out = PPT_INPUT
except PermissionError:
    prs.save(PPT_OUTPUT)
    out = PPT_OUTPUT
print(f'WROTE: {out}')
print(f'  슬라이드 {len(prs.slides)}장')
