"""사용자가 수정한 PPT에 신규 슬라이드 2장 추가.
- 슬라이드 10: 🎬 강의주소 — 신규 영역
- 슬라이드 11: 📁 그림 파일 업로드 UX

기존 슬라이드는 건드리지 않음.
"""
import os, sys, shutil
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = lambda f: os.path.join('data/presentation/img', f)
PPT = 'data/presentation/엑셀-등록페이지_연동.pptx'
BACKUP = 'data/presentation/엑셀-등록페이지_연동_백업.pptx'

# 백업
if not os.path.exists(BACKUP):
    shutil.copy(PPT, BACKUP)
    print(f'BACKUP: {BACKUP}')

FONT = '맑은 고딕'
MONO = 'Consolas'

C_PRIMARY = RGBColor(0xC2, 0x18, 0x5B)
C_DARK = RGBColor(0x1F, 0x29, 0x37)
C_MUTED = RGBColor(0x6B, 0x72, 0x80)
C_BG = RGBColor(0xFA, 0xFA, 0xFA)
C_OK = RGBColor(0x05, 0x96, 0x69)
C_WARN = RGBColor(0xD9, 0x77, 0x06)

C_VIDEO = RGBColor(0x7C, 0x3A, 0xED)   # violet
C_IMG = RGBColor(0xD9, 0x77, 0x06)     # amber
C_STEM = RGBColor(0xDB, 0x27, 0x77)
C_META = RGBColor(0x47, 0x55, 0x69)


def tb(slide, l, t, w, h, lines, *, size=14, bold=False, color=C_DARK,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return box


def rect(slide, l, t, w, h, fill, line=None, line_w=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def title_bar(slide, prs, t, sub=None):
    rect(slide, 0, 0, prs.slide_width, Inches(0.15), C_PRIMARY)
    tb(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
       t, size=26, bold=True)
    if sub:
        tb(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4),
           sub, size=13, color=C_MUTED)


prs = Presentation(PPT)
BLANK = prs.slide_layouts[6]
print(f'기존 슬라이드 {len(prs.slides)}장 — 신규 2장 append')


# =============================================================
# Slide 10 — 🎬 강의주소
# =============================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, prs, '7. 🎬 강의주소 — 신규 영역 (메타 6번째)',
          '시험 풀이 화면 상단에 강의 버튼을 띄우기 위한 URL 컬럼')

# 좌: 등록 페이지 강의 fieldset 캡처
tb(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(0.4),
   '🖊 등록 페이지 — 강의주소 입력', size=14, bold=True, color=C_VIDEO)
s.shapes.add_picture(IMG('register_video_field.png'),
                     Inches(0.4), Inches(1.9), width=Inches(6.0))

# 중앙: 21번 미리보기 상단 강의 버튼 부분
tb(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(0.4),
   '👁 미리보기 — 시험 페이지 상단 강의 버튼 시뮬레이션',
   size=14, bold=True, color=C_VIDEO)
s.shapes.add_picture(IMG('register_q21_preview.png'),
                     Inches(6.9), Inches(1.9), width=Inches(6.0))

# 하단: 영역 설명 카드
rect(s, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.95),
     RGBColor(0xED, 0xE9, 0xFE))
tb(s, Inches(0.55), Inches(6.5), Inches(0.5), Inches(0.4),
   '🎬', size=20, color=C_VIDEO)
tb(s, Inches(1.0), Inches(6.5), Inches(12), Inches(0.35),
   '강의주소 (1컬럼) — 메타 영역에 추가된 신규 컬럼',
   size=13, bold=True, color=C_VIDEO)
tb(s, Inches(1.0), Inches(6.85), Inches(12), Inches(0.5),
   [
       '· 엑셀 v3 양식: "코드" 영역 다음에 보라색 "강의" 영역으로 분리 (1컬럼)',
       '· 시험 페이지가 이 URL을 받아 상단에 "🎬 이 문제 강의 보기" 버튼 노출 — YouTube 등 외부 영상 연결',
   ], size=11)


# =============================================================
# Slide 11 — 📁 그림 파일 업로드 UX
# =============================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, prs, '8. 📁 그림 파일 업로드 UX — 모든 그림 컬럼에 적용',
          '드래그앤드롭 / 클릭 업로드 / 실제 파일 저장 + 썸네일 미리보기')

# 상단 좌: 빈 드롭존
tb(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(0.4),
   '① 빈 상태 — 드롭존 (클릭 / 드래그앤드롭)', size=13, bold=True, color=C_IMG)
s.shapes.add_picture(IMG('register_image_dropzone.png'),
                     Inches(0.4), Inches(1.9), width=Inches(6.0))

# 상단 우: 업로드 후 (해설 블록 영역에 그림 들어간 상태)
tb(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(0.4),
   '② 업로드 후 — 블록 시퀀스 안의 그림 블록',
   size=13, bold=True, color=C_OK)
s.shapes.add_picture(IMG('register_q21_explain_blocks.png'),
                     Inches(6.9), Inches(1.9), width=Inches(6.0))

# 하단: 동작 설명
rect(s, Inches(0.4), Inches(5.4), Inches(12.5), Inches(2.0),
     RGBColor(0xFE, 0xF3, 0xC7), line=C_IMG, line_w=2)

tb(s, Inches(0.55), Inches(5.5), Inches(12), Inches(0.4),
   '⚙ 동작 흐름', size=14, bold=True, color=C_IMG)

# 5단계 가로 흐름
steps = [
    ('1. 파일 선택', '클릭 또는\n드롭'),
    ('2. POST', '/api/upload\nmultipart'),
    ('3. 서버 저장', 'public/questions/\n<code>__<slot>__<name>'),
    ('4. 파일명 회신', 'JSON\n{ok, filename}'),
    ('5. 미리보기', '컬럼 값 채움\n+ 썸네일 표시'),
]
step_y = Inches(6.0)
step_w = Inches(2.3)
gap = Inches(0.1)
for i, (t, body) in enumerate(steps):
    cx = Inches(0.6 + i * (2.3 + 0.1))
    rect(s, cx, step_y, step_w, Inches(1.25), RGBColor(0xFF, 0xFF, 0xFF),
         line=C_IMG)
    tb(s, cx + Inches(0.1), step_y + Inches(0.1), step_w - Inches(0.2),
       Inches(0.35), t, size=11, bold=True, color=C_IMG)
    tb(s, cx + Inches(0.1), step_y + Inches(0.5), step_w - Inches(0.2),
       Inches(0.7), body, size=10, color=C_DARK, font=MONO)
    if i < len(steps) - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  cx + step_w + Inches(0.01),
                                  step_y + Inches(0.5),
                                  Inches(0.08), Inches(0.25))
        arr.fill.solid(); arr.fill.fore_color.rgb = C_MUTED
        arr.line.fill.background()

# 우측 하단 적용 범위
tb(s, Inches(0.55), Inches(7.05), Inches(12), Inches(0.3),
   '적용 컬럼: 발문그림 · 보기1~4그림 · 해설(블록 시퀀스 안의 그림) · 학습포인트(블록 시퀀스 안의 그림)',
   size=10, color=C_MUTED)


prs.save(PPT)
print(f'WROTE: {PPT}')
print(f'  슬라이드 {len(prs.slides)}장 (기존 9 + 신규 2)')
