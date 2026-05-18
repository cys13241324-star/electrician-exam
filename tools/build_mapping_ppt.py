"""매핑 명세 PPT — 개발자가 보고 파서를 구현할 수 있게.
PDF 영역 ↔ 엑셀 컬럼 매핑 시각화 + 추출 규칙 + 정규식 + 엣지 케이스.
출력: data/presentation/PDF-엑셀_매핑명세.pptx
"""
import os, sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = lambda f: os.path.join('data/presentation/img', f)
OUT = 'data/presentation/PDF-엑셀_매핑명세.pptx'
os.makedirs(os.path.dirname(OUT), exist_ok=True)

FONT_KR = '맑은 고딕'
FONT_MONO = 'Consolas'

# 영역별 컬러 (오버레이와 일치)
C_CAT = RGBColor(0x60, 0x7D, 0x8B)
C_NUM = RGBColor(0x45, 0x5A, 0x64)
C_STEM = RGBColor(0xE9, 0x1E, 0x63)
C_IMG = RGBColor(0xFF, 0xC1, 0x07)
C_OPT = RGBColor(0x21, 0x96, 0xF3)
C_EXP = RGBColor(0xFF, 0x6F, 0x00)
C_LP = RGBColor(0x2E, 0x7D, 0x32)
C_WR = RGBColor(0x7B, 0x1F, 0xA2)
C_ANS = RGBColor(0x00, 0x89, 0x7B)

C_DARK = RGBColor(0x21, 0x21, 0x21)
C_MUTED = RGBColor(0x75, 0x75, 0x75)
C_BG = RGBColor(0xFA, 0xFA, 0xFA)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def textbox(slide, left, top, w, h, lines, *,
            size=14, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, font=FONT_KR):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def rect(slide, left, top, w, h, fill, line=None, line_w=1):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def title(slide, t, sub=None):
    rect(slide, 0, 0, prs.slide_width, Inches(0.15), C_STEM)
    textbox(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
            t, size=26, bold=True)
    if sub:
        textbox(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4),
                sub, size=13, color=C_MUTED)


def code_box(slide, left, top, w, h, code_lines):
    rect(slide, left, top, w, h, C_CODE_BG)
    textbox(slide, left + Inches(0.1), top + Inches(0.08), w - Inches(0.2), h - Inches(0.15),
            code_lines, size=11, font=FONT_MONO, color=C_DARK)


def color_chip(slide, left, top, color, label, *, w=Inches(2.4)):
    """좌측 컬러칩 + 라벨 (영역별 색 표시)"""
    rect(slide, left, top, Inches(0.25), Inches(0.32), color)
    textbox(slide, left + Inches(0.35), top - Inches(0.02), w, Inches(0.4),
            label, size=13, bold=True)


def table_basic(slide, left, top, w, h, data, *, header_color=C_STEM, font_size=11):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, left, top, w, h).table
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = FONT_KR
                    run.font.size = Pt(font_size)
                    run.font.bold = (r == 0)
                    run.font.color.rgb = C_DARK if r > 0 else C_WHITE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
    return tbl


# =============================================================
# Slide 1 — 표지
# =============================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, prs.slide_width, prs.slide_height, C_BG)
rect(s, Inches(0.6), Inches(2.8), Inches(0.5), Inches(0.06), C_STEM)
textbox(s, Inches(0.6), Inches(2.95), Inches(12), Inches(1.2),
        'PDF ↔ 엑셀 컬럼 매핑 명세서',
        size=42, bold=True)
textbox(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.6),
        '개발자용 — 파서 / 임포터 구현 가이드',
        size=18, color=C_MUTED)
textbox(s, Inches(0.6), Inches(5.0), Inches(12), Inches(1.5),
        [
            '· 어떤 PDF 영역이 어떤 엑셀 컬럼으로 매핑되는가',
            '· 각 영역의 시작·끝 경계 식별 규칙',
            '· 텍스트 추출 후 후처리 / 정규식',
            '· 엣지 케이스 (그림 / 표 / 빈칸문제 / 수식 손상)',
        ], size=14)
textbox(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
        '기준 데이터: 2022년 1회 (60문항)  ·  v2 양식 (17 콘텐츠 컬럼)',
        size=12, color=C_MUTED)


# =============================================================
# Slide 2 — 매핑 시각화 한눈에 (페이지 2 콘덴서)
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '1. 매핑 한눈에 — 01번 콘덴서 (전체 영역 정상 케이스)',
      'PDF 한 문항 = 엑셀 한 행. 컬러별로 어느 영역이 어느 컬럼인지')

# 좌측: 컬러 매핑 이미지
s.shapes.add_picture(IMG('mapping_p02_q01.png'),
                     Inches(0.3), Inches(1.4), height=Inches(5.9))

# 우측: 컬러 칩 범례 + 엑셀 컬럼 매핑
right_x = Inches(8.3)
y = Inches(1.5)
textbox(s, right_x, y, Inches(4.7), Inches(0.5),
        '엑셀 컬럼 매핑', size=15, bold=True, color=C_STEM)
y = Inches(2.0)
mapping_chips = [
    (C_CAT,  '카테고리          ', '→ 과목ID (theory/machinery/facility)'),
    (C_NUM,  '문항번호          ', '→ 번호'),
    (C_STEM, '발문              ', '→ 발문 (HTML)'),
    (C_OPT,  '보기①~④          ', '→ 보기1·2·3·4 (HTML)'),
    (C_EXP,  '해설              ', '→ 해설 (HTML)'),
    (C_LP,   '학습 POINT 박스   ', '→ 학습포인트 (HTML, <table>)'),
    (C_WR,   '오답분석          ', '→ 오답분석 (HTML, <ol>)'),
    (C_ANS,  '정답 (➊➋➌➍)      ', '→ 정답(1~4) — 숫자'),
]
for color, src_label, dst_label in mapping_chips:
    color_chip(s, right_x, y, color, src_label)
    textbox(s, right_x + Inches(2.4), y, Inches(2.7), Inches(0.4),
            dst_label, size=12, color=C_DARK)
    y += Inches(0.45)


# =============================================================
# Slide 3 — 매핑 시각화 (페이지 16 그림 케이스)
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '2. 매핑 한눈에 — 51번 그림기호 (그림 필요 케이스)',
      '발문그림 컬럼이 있어야 하는 케이스 — 자동 감지 + 수동 마커')

s.shapes.add_picture(IMG('mapping_p16_q51.png'),
                     Inches(0.3), Inches(1.4), width=Inches(8.0))

# 우측 설명
right_x = Inches(8.5)
y = Inches(1.5)
textbox(s, right_x, y, Inches(4.7), Inches(0.5),
        '그림 영역 처리', size=15, bold=True, color=C_IMG)
y = Inches(2.0)
rect(s, right_x, y, Inches(4.7), Inches(4.3), C_BG)
textbox(s, right_x + Inches(0.15), y + Inches(0.1), Inches(4.5), Inches(4.1),
        [
            '① 파서 자동 감지',
            '   발문/해설에 "그림과 같은",',
            '   "아래 그림", "그림기호",',
            '   "심볼" 키워드 발견 시',
            '   → [필요: …] 마커 자동 삽입',
            '',
            '② 엑셀에 표기되는 형태',
            '   발문그림 = "[필요: 기호심볼',
            '              — 수동조작 접점]"',
            '',
            '③ 운영자 후속 작업',
            '   • 그림 파일 업로드',
            '   • 컬럼 값 → 실제 파일명',
            '     (예: q51_manual_contact.png)',
            '',
            '※ PDF 텍스트 추출 단계에서',
            '   그림 자체는 변환 불가',
        ], size=12)


# =============================================================
# Slide 4 — 텍스트 추출 경계 규칙 (전체 흐름)
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '3. 추출 경계 규칙 — 한 문항 안에서 영역을 어떻게 자르는가',
      '문항번호 → 발문 → 보기 → 해설(+학습포인트) → 오답분석 → 끝 태그 순')

# 흐름도: 세로 스택
stages = [
    ('1. 문항 시작', '문항번호 단독 라인  /^\\d{2}$/', '"01", "02", …, "60"', C_NUM),
    ('2. 발문', '문항번호 직후 ~ 첫 ① 직전', '발문 텍스트', C_STEM),
    ('3. 발문그림', '발문에 "그림" 키워드 + 발문~보기 사이 공백 영역', '자동 마커 삽입', C_IMG),
    ('4. 보기 ①~④', '①②③④ 마커 — 단, ( ① ) 형태는 빈칸으로 별도 처리', '보기1~4 텍스트', C_OPT),
    ('5. 해설', '보기 ④ 끝 + 공백라인 이후 ~ "오답분석" 또는 "공식/법칙" 헤더 직전', '해설 텍스트', C_EXP),
    ('6. 학습포인트', '해설 안의 "OO 공식/법칙/특징/분류/관계" 헤더 라인 ~ 끝 태그 직전', '학습포인트 (HTML 표 포함)', C_LP),
    ('7. 오답분석', '"오답분석" 라인 ~ 끝 태그 직전', '오답분석 (HTML <ol>)', C_WR),
    ('8. 정답', '끝 태그 /^\\[전기이론\\]\\s\\(00p\\)\\s([➊➋➌➍])$/', '➊➋➌➍ → 1·2·3·4', C_ANS),
]
yy = Inches(1.5)
row_h = Inches(0.62)
for i, (step, rule, result, color) in enumerate(stages):
    y = yy + i * row_h
    # 좌측 번호 박스
    rect(s, Inches(0.4), y, Inches(2.0), Inches(0.55), color)
    textbox(s, Inches(0.4), y + Inches(0.07), Inches(2.0), Inches(0.4),
            step, size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # 중앙 규칙
    textbox(s, Inches(2.6), y + Inches(0.05), Inches(7.5), Inches(0.5),
            rule, size=11, color=C_DARK, font=FONT_MONO)
    # 우측 결과
    textbox(s, Inches(10.3), y + Inches(0.05), Inches(2.8), Inches(0.5),
            '→ ' + result, size=11, color=color, bold=True)


# =============================================================
# Slide 5 — 발문 추출 상세
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '4. 발문 컬럼 — 추출 규칙', '문항번호 직후부터 첫 ① 마커 직전까지')

# 좌측: 입력 예시
rect(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(5.5), C_CODE_BG)
textbox(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.4),
        '입력 (pypdf 추출)', size=13, bold=True, color=C_STEM)
textbox(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0),
        [
            '01',
            '콘덴서의 정전용량에 대한 설명으로 틀린 것은?',
            '① 전압에 반비례한다.',
            '② 이동 전하량에 비례한다.',
            '③ 극판의 넓이에 비례한다.',
            '④ 극판의 간격에 비례한다.',
            '',
            '— 발문 영역 —',
            ' 시작: "01" 다음 줄',
            ' 끝:   첫 "①" 이전',
        ], size=12, font=FONT_MONO)

# 우측: 규칙
right_x = Inches(6.7)
textbox(s, right_x, Inches(1.5), Inches(6.3), Inches(0.4),
        '규칙', size=15, bold=True, color=C_STEM)
textbox(s, right_x, Inches(2.0), Inches(6.3), Inches(5.0),
        [
            '◆ 시작 경계',
            '   문항번호 패턴 /^\\s*(\\d{2})\\s*$/ 매치한 다음 라인',
            '',
            '◆ 끝 경계',
            '   ① 마커 첫 출현. 단:',
            '   • ( ① ) 형태 (괄호 안) → 빈칸 표시, 발문에 포함',
            '   • OPT_INLINE 정규식:',
            '     /(?<![가-힣A-Za-z0-9])([①②③④])/',
            '',
            '◆ 후처리',
            '   • 자모 분리 패턴 "기 본 파" → "기본파"',
            '     /([가-힣]\\s){2,}[가-힣]/ 매치 영역만 공백 제거',
            '   • PUA 영역 (U+E000–U+E0FF) → "[수식]" 마커',
            '   • 단위 공백 정리: "kV A" → "kVA"',
            '',
            '◆ 출력 컬럼: 발문 (HTML)',
            '   허용 태그: <br>, <sub>, <sup>, <strong>, <i>, <em>',
            '',
            '◆ 자동 그림 감지 (발문 → 발문그림 컬럼)',
            '   /그림과\\s*같은|아래\\s*그림|그림기호|심볼/ 매치 시',
            '   발문그림 = "[필요: 회로도/심볼/도면]"',
        ], size=11)


# =============================================================
# Slide 6 — 보기 추출 상세
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '5. 보기①~④ 컬럼 — 마커 기반 분할',
      '①②③④ 마커 위치로 4구간 자르기. 빈칸문제는 예외 처리.')

# 좌측: 정상 케이스
rect(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(2.6), C_CODE_BG)
textbox(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.4),
        '정상 케이스', size=13, bold=True, color=C_OPT)
textbox(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.0),
        [
            '① 전압에 반비례한다.',
            '② 이동 전하량에 비례한다.',
            '③ 극판의 넓이에 비례한다.',
            '④ 극판의 간격에 비례한다.',
            '',
            '→ 보기1="전압에 반비례한다." 등',
        ], size=12, font=FONT_MONO)

# 좌측 하단: 빈칸문제
rect(s, Inches(0.4), Inches(4.3), Inches(6.0), Inches(2.8), RGBColor(0xFF, 0xF0, 0xF0))
textbox(s, Inches(0.5), Inches(4.4), Inches(5.8), Inches(0.4),
        '예외 — 빈칸문제 (9번)', size=13, bold=True, color=C_WR)
textbox(s, Inches(0.5), Inches(4.8), Inches(5.8), Inches(2.2),
        [
            '발문: "...코일을 지나는 ( ① )의 매초',
            '       변화량과 코일의 ( ② )에 비례한다."',
            '',
            '보기:',
            '  ①① 자속, ② 굵기   ②① 자속, ② 권수',
            '  ③① 전류, ② 권수   ④① 전류, ② 굵기',
            '',
            '→ 발문 안 (①)(②)와 보기 ①②③④ 충돌',
        ], size=11, font=FONT_MONO)

# 우측: 규칙
right_x = Inches(6.7)
textbox(s, right_x, Inches(1.5), Inches(6.3), Inches(0.4),
        '규칙', size=15, bold=True, color=C_STEM)
textbox(s, right_x, Inches(2.0), Inches(6.3), Inches(5.3),
        [
            '◆ 분할 마커 (OPT_INLINE)',
            '   /(?<![가-힣A-Za-z0-9])([①②③④])/',
            '   → 앞 글자가 한글/영문/숫자가 아닐 때만 마커',
            '',
            '◆ 빈칸 패턴 전처리 (BLANK_FILL_RE)',
            '   /\\(\\s*([①②③④])\\s*\\)/',
            '   "( ① )" → "(빈칸1)"  보기 마커와 분리',
            '   이슈 컬럼에 "빈칸문제" 플래그 기록',
            '',
            '◆ 분할 알고리즘',
            '   1. OPT_INLINE으로 모든 마커 매치 위치 수집',
            '   2. 처음 보는 마커 순서대로 ①②③④ 할당',
            '   3. 각 마커 시작 ~ 다음 마커 시작 사이 텍스트 추출',
            '',
            '◆ 그림형 보기 처리',
            '   → 보기 텍스트 비어 있으면 보기N그림 컬럼 사용',
            '   → 자동 변환 불가 → 수동 마커',
            '',
            '◆ 출력 컬럼: 보기1, 보기2, 보기3, 보기4 (HTML)',
            '   허용 태그: <i>, <sub>, <sup>, √, ε, π 등 HTML 엔티티',
        ], size=11)


# =============================================================
# Slide 7 — 정답 추출 + 인코딩 변환표
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '6. 정답(1~4) 컬럼 — 끝 태그에서 추출 + 유니코드 변환',
      '문항 끝의 ➊➋➌➍ → 1·2·3·4')

# 좌측: 끝 태그 형식
rect(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(3.5), C_CODE_BG)
textbox(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.4),
        '끝 태그 형식', size=13, bold=True, color=C_ANS)
textbox(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.0),
        [
            '[전기이론] (00p) ➍',
            '[전기이론] (00p) ➊',
            '[전기이론] (00p) ➋',
            '[전기이론] (00p) ➌',
            '',
            'END_TAG_LINE 정규식:',
            '  /^\\[(전기이론|전기기기|전기설비)\\]',
            '    \\s\\(\\d+p\\)\\s([➊➋➌➍])$/',
            '',
            '※ 카테고리는 신뢰 안 함 (PDF 오류로 전부 [전기이론])',
            '※ 정답 마커만 캡처',
        ], size=12, font=FONT_MONO)

# 우측: 인코딩 변환표
right_x = Inches(6.7)
textbox(s, right_x, Inches(1.5), Inches(6.3), Inches(0.4),
        '인코딩 변환표', size=15, bold=True, color=C_STEM)
table_basic(s, right_x, Inches(2.0), Inches(6.3), Inches(2.4),
            [
                ['문자', '유니코드', '의미', '엑셀 출력'],
                ['➊', 'U+278A', 'Dingbat 1', '1'],
                ['➋', 'U+278B', 'Dingbat 2', '2'],
                ['➌', 'U+278C', 'Dingbat 3', '3'],
                ['➍', 'U+278D', 'Dingbat 4', '4'],
            ], header_color=C_ANS, font_size=12)

textbox(s, right_x, Inches(4.6), Inches(6.3), Inches(0.4),
        '엣지 케이스', size=13, bold=True, color=C_WR)
textbox(s, right_x, Inches(5.0), Inches(6.3), Inches(2.0),
        [
            '· 끝 태그 자체가 누락된 문항 존재',
            '  (예: 55번 — PDF 마스터 오류)',
            '  → 정답 추출 실패 → 이슈 컬럼에 기록',
            '  → manual_fixes로 수동 보완',
            '',
            '· 출력 컬럼: 정답(1~4) — 정수형',
            '  (HTML 아님, 1·2·3·4 중 하나)',
        ], size=12)


# =============================================================
# Slide 8 — 해설 + 학습포인트 분리
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '7. 해설 / 학습포인트 — 두 컬럼으로 자동 분리',
      'PDF 한 덩어리 해설 안에 "공식 / 법칙 / 특징" 박스 → 별도 컬럼')

# 좌측: 분리 전
rect(s, Inches(0.4), Inches(1.5), Inches(5.8), Inches(5.5), C_CODE_BG)
textbox(s, Inches(0.5), Inches(1.6), Inches(5.6), Inches(0.4),
        '분리 전 (원본 해설 영역)', size=13, bold=True, color=C_EXP)
textbox(s, Inches(0.5), Inches(2.0), Inches(5.6), Inches(5.0),
        [
            '전기분해를 통해 석출된 물질의 양은',
            '통과한 전기량 및 화학당량에 각각',
            '비례한다.',
            '',
            '─────────── ← 자동 감지',
            'LP_HEADER_RE 매치 라인',
            '─────────── ',
            '',
            '패러데이의 법칙       ← 헤더',
            ' W = kIt [g]',
            ' • W: 석출된 물질의 양[g]',
            ' • k: 전기 화학 당량[g/C]',
            ' • I: 전류[A]',
            ' • t: 시간[s]',
        ], size=11, font=FONT_MONO)

# 우측: 규칙 + 분리 후
right_x = Inches(6.5)
textbox(s, right_x, Inches(1.5), Inches(6.5), Inches(0.4),
        '학습포인트 헤더 패턴', size=15, bold=True, color=C_LP)
code_box(s, right_x, Inches(2.0), Inches(6.5), Inches(1.4),
         [
             'LP_HEADER_RE = /^.{2,30}(공식|법칙|특징|',
             '  분류|종류|성질|관계|표현식|효율|구분|',
             '  반작용|기준|구조|길이|깊이|주파수|표준|',
             '  화학반응|순시값)\\s*$/',
             '',
             '→ 짧은 라인이 위 키워드로 끝나면 분리 트리거',
         ])

textbox(s, right_x, Inches(3.6), Inches(6.5), Inches(0.4),
        '분리 후 — 두 컬럼', size=13, bold=True, color=C_STEM)
table_basic(s, right_x, Inches(4.0), Inches(6.5), Inches(3.0),
            [
                ['컬럼', '내용'],
                ['해설', '"전기분해를 통해 석출된 물질의 양은 통과한 전기량 및 화학당량에 각각 비례한다."'],
                ['학습포인트', '<p><strong>패러데이의 법칙</strong></p><p><i>W</i> = <i>kIt</i> [g]</p><ul><li>W: 석출량[g]</li><li>k: 화학당량[g/C]</li><li>I: 전류[A]</li><li>t: 시간[s]</li></ul>'],
            ], header_color=C_LP, font_size=10)


# =============================================================
# Slide 9 — 오답분석 추출
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '8. 오답분석 컬럼 — "오답분석" 키워드 트리거',
      '해설 영역과 분리. <ol> 리스트로 HTML 변환 권장.')

# 좌측: 원본
rect(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(5.5), C_CODE_BG)
textbox(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.4),
        '원본 (pypdf 추출)', size=13, bold=True, color=C_WR)
textbox(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0),
        [
            '... 해설 본문 ...',
            'TRIAC이다.',
            '오답분석            ← 시작 트리거',
            '①, ② 발진 회로 등에 사용되는 소자이다.',
            '③ 양방향성 2단자 소자로, 주로 TRIAC의',
            '   트리거용으로 사용된다.',
            '[전기이론] (00p) ➍    ← 끝 (정답 태그)',
        ], size=12, font=FONT_MONO)

# 우측
right_x = Inches(6.7)
textbox(s, right_x, Inches(1.5), Inches(6.3), Inches(0.4),
        '추출 규칙', size=15, bold=True, color=C_STEM)
textbox(s, right_x, Inches(2.0), Inches(6.3), Inches(2.5),
        [
            '◆ 시작 경계',
            '   line.strip().startswith("오답분석")',
            '',
            '◆ 끝 경계',
            '   END_TAG_LINE 매치 직전',
            '   또는 학습포인트 헤더 (LP_HEADER_RE) 등장 시',
            '',
            '◆ 출력 컬럼: 오답분석 (HTML)',
            '   권장: <ol><li>...</li></ol>',
            '   정답이 ②인 경우 li value="2"부터 시작 가능',
        ], size=12)

textbox(s, right_x, Inches(4.6), Inches(6.3), Inches(0.4),
        'HTML 변환 예시', size=13, bold=True, color=C_LP)
code_box(s, right_x, Inches(5.0), Inches(6.3), Inches(2.0),
         [
             '<ol>',
             '  <li>① ② 발진 회로 등에 사용</li>',
             '  <li value="3">양방향성 2단자 소자,',
             '       주로 TRIAC 트리거용</li>',
             '</ol>',
         ])


# =============================================================
# Slide 10 — 그림 영역 처리
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '9. 그림 컬럼 — 자동 감지 + 수동 마커',
      'PDF 텍스트 추출로는 그림 자체를 가져올 수 없음. 위치 정보만 명시.')

# 좌측: 자동 감지 규칙
rect(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(5.5), RGBColor(0xFF, 0xF8, 0xE1))
textbox(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.4),
        '자동 감지 규칙 (detect_image_need)', size=14, bold=True, color=C_IMG)

table_basic(s, Inches(0.5), Inches(2.1), Inches(5.8), Inches(2.5),
            [
                ['텍스트 패턴', '마커 출력'],
                ['"그림과 같은"', '[필요: 회로도/도면]'],
                ['"그림기호"', '[필요: 기호 심볼]'],
                ['"아래 그림"', '[필요: 도면/그림]'],
                ['"다음 그림"', '[필요: 도면/그림]'],
                ['"심볼"', '[필요: 심볼]'],
            ], header_color=C_IMG, font_size=11)

textbox(s, Inches(0.5), Inches(4.8), Inches(5.8), Inches(0.4),
        '적용 대상 컬럼', size=13, bold=True, color=C_STEM)
textbox(s, Inches(0.5), Inches(5.2), Inches(5.8), Inches(1.8),
        [
            '· 발문에 키워드 발견 → 발문그림 컬럼',
            '· 해설에 키워드 발견 → 해설그림 컬럼',
            '· 학습포인트에 키워드 발견 → 학습포인트그림 컬럼',
            '· 보기 자체가 그림형이면 → 보기N그림 컬럼',
            '  (자동 감지 어려움 → 수동 마커 권장)',
        ], size=11)

# 우측: 수동 마커 (IMAGE_NEEDS)
right_x = Inches(6.7)
textbox(s, right_x, Inches(1.5), Inches(6.3), Inches(0.4),
        '수동 마커 (manual_fixes.IMAGE_NEEDS)', size=14, bold=True, color=C_STEM)
textbox(s, right_x, Inches(2.0), Inches(6.3), Inches(0.4),
        '자동 감지로 못 잡는 케이스 직접 명시', size=11, color=C_MUTED)

code_box(s, right_x, Inches(2.4), Inches(6.3), Inches(1.5),
         [
             'IMAGE_NEEDS = {',
             '  51: {"발문그림": "[필요: 기호심볼 ',
             '         — 수동조작 접점]"},',
             '  39: {"해설그림": "[필요: TRIAC 심볼]"},',
             '}',
         ])

textbox(s, right_x, Inches(4.1), Inches(6.3), Inches(0.4),
        '운영자 후속 작업', size=13, bold=True, color=C_LP)
textbox(s, right_x, Inches(4.5), Inches(6.3), Inches(2.5),
        [
            '1. 엑셀에서 그림 컬럼 [필요: …] 마커 검색',
            '2. PDF 원본 페이지 보면서 그림 캡처',
            '3. public/questions/ 폴더에 PNG 저장',
            '4. 엑셀 컬럼 값 → 파일명 교체',
            '   예: "q51_manual_contact.png"',
            '5. 본문 인라인 필요 시 <img src="…"> 직접 삽입',
        ], size=11)


# =============================================================
# Slide 11 — 메타데이터 + 타입 요약
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '10. 메타데이터 / 컬럼 타입 요약',
      '문항번호·카테고리·서브토픽 매핑 + 17개 컬럼 데이터 타입')

# 좌측: 메타 추출
rect(s, Inches(0.4), Inches(1.5), Inches(5.5), Inches(5.5), C_BG)
textbox(s, Inches(0.5), Inches(1.6), Inches(5.3), Inches(0.4),
        '메타데이터 추출', size=14, bold=True, color=C_NUM)

table_basic(s, Inches(0.5), Inches(2.1), Inches(5.3), Inches(2.8),
            [
                ['컬럼', '추출 위치'],
                ['번호', '문항번호 라인 (1~60 정수)'],
                ['과목ID', '페이지 안의 카테고리 헤더 라인'],
                ['', '("전기이론" → theory)'],
                ['', '("전기 기기" → machinery)'],
                ['', '("전기 설비" → facility)'],
                ['토픽ID', '커리큘럼ID 시트 참조 (수동)'],
                ['서브토픽ID', '커리큘럼ID 시트 참조 (수동)'],
            ], header_color=C_NUM, font_size=11)

textbox(s, Inches(0.5), Inches(5.1), Inches(5.3), Inches(0.4),
        '⚠️ 마스터 데이터 오류 처리', size=13, bold=True, color=C_WR)
textbox(s, Inches(0.5), Inches(5.5), Inches(5.3), Inches(1.5),
        [
            'PDF의 끝 태그 [전기이론]은 모든 문항에',
            '잘못 박혀있음 (마스터 오류).',
            '',
            '→ 카테고리는 끝 태그 무시,',
            '   페이지 헤더(전기이론/기기/설비)로 결정',
        ], size=11)

# 우측: 컬럼 타입 요약
right_x = Inches(6.1)
textbox(s, right_x, Inches(1.5), Inches(7.0), Inches(0.4),
        '17개 콘텐츠 컬럼 데이터 타입', size=14, bold=True, color=C_STEM)
table_basic(s, right_x, Inches(2.0), Inches(7.0), Inches(5.0),
            [
                ['컬럼', '타입', '필수', 'HTML', 'sanitize'],
                ['발문', '문자열', '✓', '✓', '필요'],
                ['조건', '문자열', '', '✓', '필요'],
                ['발문그림', '파일명/마커', '', '—', ''],
                ['보기1~4', '문자열', '✓', '✓', '필요'],
                ['보기1~4그림', '파일명/마커', '', '—', ''],
                ['정답(1~4)', '정수 (1·2·3·4)', '✓', '—', ''],
                ['해설', '문자열', '✓', '✓ (<table>)', '필요'],
                ['해설그림', '파일명/마커', '', '—', ''],
                ['오답분석', '문자열', '', '✓ (<ol>)', '필요'],
                ['학습포인트', '문자열', '', '✓ (<table>)', '필요'],
                ['학습포인트그림', '파일명/마커', '', '—', ''],
            ], header_color=C_STEM, font_size=10)


# =============================================================
# Slide 12 — 개발자 체크리스트 / 엣지 케이스
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '11. 개발자 체크리스트 — 구현 시 반드시 확인',
      '파서 구현 / 임포터 구현 양쪽 모두')

# 좌측: 파서 구현
rect(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(5.7), RGBColor(0xE8, 0xF0, 0xFE))
textbox(s, Inches(0.6), Inches(1.55), Inches(5.8), Inches(0.5),
        '🔧 파서 구현 (PDF → 엑셀)', size=16, bold=True, color=C_STEM)
textbox(s, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.9),
        [
            '☐ 페이지 1 = 이전 회차 마지막 문항 → 스킵',
            '☐ 페이지 헤더/푸터 정규식 노이즈 제거',
            '☐ 정답 누락 케이스 (예: 55번) — 이슈 기록',
            '☐ 빈칸문제 ( ① ) 전처리 (보기 마커와 분리)',
            '☐ PUA 영역 U+E000–U+E0FF → "[수식]" 치환',
            '☐ 자모 분리 정규식은 3글자 이상 연속만',
            '☐ 카테고리는 끝 태그 무시, 페이지 헤더 사용',
            '☐ 학습포인트 분리: LP_HEADER_RE 위치 후보',
            '☐ 그림 자동 감지 → [필요: …] 마커',
            '☐ manual_fixes로 수동 보완 — 덮어쓰기 우선순위',
            '☐ 이슈 컬럼에 케이스별 플래그 기록',
            '☐ 변환통계 시트 자동 생성 (총·이슈·정상 카운트)',
        ], size=12)

# 우측: 임포터 구현
rect(s, Inches(6.8), Inches(1.4), Inches(6.2), Inches(5.7), RGBColor(0xE8, 0xF8, 0xE8))
textbox(s, Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.5),
        '📥 임포터 구현 (엑셀 → DB/JSON)', size=16, bold=True, color=C_LP)
textbox(s, Inches(7.0), Inches(2.1), Inches(5.8), Inches(4.9),
        [
            '☐ HTML sanitize (DOMPurify 등으로 XSS 차단)',
            '   허용 화이트리스트: <i><sub><sup><strong><em>',
            '   <br><p><ul><ol><li><table><thead><tbody>',
            '   <tr><th><td>',
            '   금지: <script><iframe><style> + 모든 on* 속성',
            '☐ 그림 컬럼 값이 "[필요: …]" 마커면 업로드 차단',
            '   또는 별도 검수 큐로 이동',
            '☐ 정답(1~4) 정수 검증 (1~4 범위 외 거부)',
            '☐ 필수 컬럼 누락 시 행 거부',
            '   (발문/보기1~4/정답/해설)',
            '☐ <em>원본 확인 필요</em> 태그 있으면 검수 큐',
            '☐ 모바일·태블릿 렌더링 검증',
            '   특히 <table>이 좁은 화면에서 깨지지 않는지',
            '☐ 보기N그림과 보기N 텍스트 둘 다 채워졌을 때',
            '   화면 레이아웃 정책 (그림 우선? 캡션 표시?)',
        ], size=11)

# 저장
prs.save(OUT)
print(f'WROTE: {OUT}')
print(f'  슬라이드 {len(prs.slides)}장')
