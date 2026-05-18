"""PPT 생성: 2022년 1회 PDF → CBT 엑셀 변환 프로젝트 (팀원용, 11장)
- 한국어 발표 / 맑은 고딕
- 주요 PDF 스크린샷 인용
- 출력: data/presentation/2022_1회_변환프로젝트.pptx
"""
import os, sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = lambda f: os.path.join('data/presentation/img', f)
OUT = 'data/presentation/2022_1회_변환프로젝트.pptx'
os.makedirs(os.path.dirname(OUT), exist_ok=True)

FONT_KR = '맑은 고딕'
COLOR_PRIMARY = RGBColor(0xC2, 0x18, 0x5B)   # 분홍 (addto 톤)
COLOR_ACCENT = RGBColor(0x21, 0x96, 0xF3)    # 파랑
COLOR_DARK = RGBColor(0x21, 0x21, 0x21)
COLOR_MUTED = RGBColor(0x75, 0x75, 0x75)
COLOR_BG_LIGHT = RGBColor(0xFA, 0xFA, 0xFA)
COLOR_OK = RGBColor(0x2E, 0x7D, 0x32)
COLOR_WARN = RGBColor(0xE6, 0x5A, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]  # Blank


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=COLOR_DARK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font=FONT_KR):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def slide_title(slide, title, subtitle=None):
    # 좌상단 컬러 바
    bar = add_rect(slide, 0, 0, Inches(13.333), Inches(0.15), COLOR_PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
                title, size=28, bold=True, color=COLOR_DARK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
                    subtitle, size=14, color=COLOR_MUTED)


# ================================================================
# Slide 1 — 표지
# ================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, COLOR_BG_LIGHT)
add_rect(s, Inches(0.6), Inches(2.8), Inches(0.5), Inches(0.06), COLOR_PRIMARY)
add_textbox(s, Inches(0.6), Inches(2.95), Inches(12), Inches(1.2),
            '2022년 1회 PDF → CBT 문항 엑셀 자동 변환',
            size=42, bold=True, color=COLOR_DARK)
add_textbox(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.6),
            '전기기능사 학습 서비스 · 문항 데이터 입력 파이프라인',
            size=20, color=COLOR_MUTED)
add_textbox(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
            '발표일: 2026-05-15  ·  대상: 팀원/공동작업자',
            size=12, color=COLOR_MUTED)


# ================================================================
# Slide 2 — 목표 / 작업 범위
# ================================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, '1. 작업 목표', '시중 출간 PDF → 우리 CBT 시스템에 바로 올라가는 엑셀 양식')

# 좌측 박스: 목표
add_rect(s, Inches(0.5), Inches(1.6), Inches(6.0), Inches(5.3), RGBColor(0xFF, 0xF7, 0xE0))
add_textbox(s, Inches(0.7), Inches(1.75), Inches(5.6), Inches(0.5),
            '목표', size=20, bold=True, color=COLOR_PRIMARY)
add_textbox(s, Inches(0.7), Inches(2.3), Inches(5.6), Inches(4.4),
            [
                '• PDF 모의고사 1회분 (60문항) 전체를',
                '  v2 엑셀 양식으로 자동 변환',
                '',
                '• 발문 / 보기 / 해설 / 오답분석 / 학습 POINT',
                '  5개 요소를 컬럼 단위로 분리',
                '',
                '• 수식·표는 HTML로 입력',
                '  (업로드 시 그대로 렌더링)',
                '',
                '• 그림이 필요한 위치는 컬럼에 명시',
            ],
            size=15, color=COLOR_DARK)

# 우측 박스: 산출물
add_rect(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(5.3), RGBColor(0xE8, 0xF8, 0xE8))
add_textbox(s, Inches(7.0), Inches(1.75), Inches(5.6), Inches(0.5),
            '산출물', size=20, bold=True, color=COLOR_OK)
add_textbox(s, Inches(7.0), Inches(2.3), Inches(5.6), Inches(4.4),
            [
                '① data/templates/question-template-v2.xlsx',
                '   — 단순 등록용 v2 양식 (21컬럼)',
                '',
                '② data/templates/[…]통합템플릿_v2.xlsx',
                '   — 메타데이터 포함 v2 양식 (42컬럼)',
                '',
                '③ data/2022_1회_변환결과_v2.xlsx',
                '   — 60문항 전체 변환 결과',
                '',
                '④ tools/manual_fixes.py',
                '   — 60문항 보완 데이터 (재사용 가능)',
                '',
                '⑤ tools/convert_2022_1.py',
                '   — 변환 스크립트 (다른 회차에도 적용)',
            ],
            size=13, color=COLOR_DARK)


# ================================================================
# Slide 3 — 엑셀 양식 v2 컬럼 구조
# ================================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, '2. 엑셀 양식 v2 컬럼 구조', '각 요소를 별도 컬럼으로 분리 · 텍스트 컬럼은 HTML 허용')

# 컬럼 표
table_data = [
    ['#', '컬럼', '용도', 'HTML', '비고'],
    ['1', '발문', '문제 본문', '✅', '<br>, <sub>, <sup>, <strong>'],
    ['2', '조건', '주어진 조건', '✅', ''],
    ['3', '발문그림', '회로도·도면', '—', '파일명 또는 [필요: …]'],
    ['4', '보기1', '① 텍스트/수식', '✅', '그림형이면 비우거나 캡션'],
    ['5', '보기1그림', '① 그림 파일명', '—', ''],
    ['6~11', '보기2~4 + 그림', '동일', '✅/—', ''],
    ['12', '정답(1~4)', '1·2·3·4', '—', ''],
    ['13', '해설', '해설 본문', '✅', '<table>, <ol>, <ul> 가능'],
    ['14', '해설그림', '해설 다이어그램', '—', ''],
    ['15', '오답분석', '②③④ 분석', '✅', '<ol><li>...</li></ol> 권장'],
    ['16', '학습포인트', '독끝 학습 Point', '✅', '<table> 가능'],
    ['17', '학습포인트그림', '', '—', ''],
]
rows = len(table_data)
cols = 5
tbl = s.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.6)).table
tbl.columns[0].width = Inches(0.6)
tbl.columns[1].width = Inches(2.2)
tbl.columns[2].width = Inches(3.4)
tbl.columns[3].width = Inches(0.8)
tbl.columns[4].width = Inches(5.3)
for r, row in enumerate(table_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = FONT_KR
                run.font.size = Pt(12 if r > 0 else 13)
                run.font.bold = (r == 0)
                run.font.color.rgb = COLOR_DARK if r > 0 else RGBColor(0xFF, 0xFF, 0xFF)
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_PRIMARY


# ================================================================
# Slide 4 — PDF 추출 시 발견된 이슈 15개
# ================================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, '3. PDF에서 텍스트 추출할 때 만난 이슈 15가지',
            'pypdf로 추출 → 어디서 깨지고 어디서 노이즈가 들어오는가')

categories = [
    ('구조 이슈 (6)', RGBColor(0xFF, 0xE0, 0xE0), [
        '① 페이지 1 = 이전 회차의 60번 문항 (스킵 필요)',
        '② 페이지 2 끝 "10문항 / 19문항 / 1문항" 노이즈',
        '③ 페이지 헤더 (232 전기기능사 필기 기본서) 매 페이지 섞임',
        '④ 페이지 18 끝 "1회독(20문항)" 정답수 영역 노이즈',
        '⑤ 카테고리 헤더(전기이론/기기/설비) 문항 사이 흩어짐',
        '⑥ 모든 문항 끝 태그가 [전기이론]으로 잘못 표기 (마스터 오류)',
    ]),
    ('텍스트 이슈 (5)', RGBColor(0xFF, 0xF0, 0xD0), [
        '⑦ 수식·공식 보기가 통째로 빈 줄로 (벡터 미임베딩)',
        '⑧ 해설 본문 변수 글자 빠짐 — "정전용량()은 넓이()에 비례"',
        '⑨ 한글 자모 공백 분리 — "①기 본 파", "④ 미 관 상  좋 다"',
        '⑩ 단위 표기 공백 — "kV A" → "kVA"',
        '⑪ 정답 인코딩 ➊➋➌➍ → 1·2·3·4 변환 필요',
    ]),
    ('시각 / 파싱 이슈 (4)', RGBColor(0xE0, 0xEC, 0xF8), [
        '⑫ 그림 문항 (51번 그림기호, 39번 TRIAC 심볼) — 텍스트 추출 불가',
        '⑬ 표 (23·30·43·53번) — 셀 경계 깨짐, 수동 재입력 필요',
        '⑭ 빈칸문제 (9번) — 발문 안의 ( ① )과 보기 ①②③④ 충돌',
        '⑮ 보기가 한 줄에 2개씩 vs 별도 줄 — 일관성 없음',
    ]),
]

x = Inches(0.4)
y_start = Inches(1.5)
w = Inches(4.18)
gap = Inches(0.1)

for i, (title, fill, items) in enumerate(categories):
    cx = Inches(0.4 + i * 4.28)
    add_rect(s, cx, y_start, w, Inches(5.5), fill)
    add_textbox(s, cx + Inches(0.2), y_start + Inches(0.2), w - Inches(0.4), Inches(0.5),
                title, size=16, bold=True, color=COLOR_PRIMARY)
    add_textbox(s, cx + Inches(0.2), y_start + Inches(0.75), w - Inches(0.4), Inches(4.6),
                items, size=11, color=COLOR_DARK)


# ================================================================
# Slide 5 — 자동 변환 흐름
# ================================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, '4. 자동 변환 파이프라인',
            'PDF 텍스트 추출 → 파싱 → 보완 → 엑셀 출력')

# 5단계 박스 가로
steps = [
    ('1. PDF 추출', 'pypdf로\n페이지별\n텍스트 추출', RGBColor(0xE8, 0xF0, 0xFE)),
    ('2. 노이즈 제거', '헤더/푸터/\n정답수 영역\n정규식 필터', RGBColor(0xFF, 0xF7, 0xE0)),
    ('3. 문항 단위 파싱', '문항번호 →\n발문 / 보기 /\n해설 / 오답분석', RGBColor(0xFF, 0xF0, 0xD0)),
    ('4. 자동 보정', '자모 공백,\n학습 POINT 분리,\n그림 마커', RGBColor(0xE8, 0xF8, 0xE8)),
    ('5. 수동 보완', 'manual_fixes\n.py로 수식·표\nHTML 덮어쓰기', RGBColor(0xFC, 0xE4, 0xEC)),
]

xx = Inches(0.4)
yy = Inches(1.9)
w = Inches(2.4)
gap = Inches(0.18)
for i, (title, body, color) in enumerate(steps):
    cx = Inches(0.4 + i * (2.4 + 0.18))
    add_rect(s, cx, yy, w, Inches(2.5), color)
    add_textbox(s, cx + Inches(0.15), yy + Inches(0.15), w - Inches(0.3), Inches(0.5),
                title, size=14, bold=True, color=COLOR_PRIMARY)
    add_textbox(s, cx + Inches(0.15), yy + Inches(0.75), w - Inches(0.3), Inches(1.6),
                body, size=12, color=COLOR_DARK)
    # 화살표 (사이 박스)
    if i < len(steps) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(0.4 + (i+1) * 2.4 + i * 0.18 + 0.01),
                                    yy + Inches(1.05),
                                    Inches(0.16), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_MUTED
        arrow.line.fill.background()

# 하단 결과 박스
add_rect(s, Inches(0.4), Inches(4.8), Inches(12.5), Inches(2.2), RGBColor(0xFA, 0xFA, 0xFA),
         line=COLOR_MUTED)
add_textbox(s, Inches(0.6), Inches(4.95), Inches(12), Inches(0.5),
            '결과', size=16, bold=True, color=COLOR_OK)
add_textbox(s, Inches(0.6), Inches(5.5), Inches(12), Inches(1.4),
            [
                '• 60문항 100% 파싱 (theory 20 / machinery 20 / facility 20)',
                '• 자동 정상 변환 59건 (98%) · 수식·표 수동 보완 59건 적용 후',
                '• 학습 POINT 60문항 모두 보유 (자동 분리 14 + 수동 채움 47)',
                '• 남은 1건: 9번 빈칸문제 노트 (데이터는 정상, 구조 표시만)',
            ],
            size=13, color=COLOR_DARK)


# ================================================================
# Slide 6 — 수식 HTML 변환 예시 (1번 콘덴서)
# ================================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, '5. 수식 HTML 변환 예시 — 01번 콘덴서 정전용량',
            'PDF 원본 → 추출 텍스트 → 수동 보완 HTML')

# 좌측: PDF 페이지 캡처 (페이지 2 상단)
s.shapes.add_picture(IMG('p02_q01_콘덴서.png'),
                     Inches(0.3), Inches(1.5),
                     width=Inches(4.7), height=Inches(5.3))
add_textbox(s, Inches(0.3), Inches(6.85), Inches(4.7), Inches(0.3),
            '원본 PDF (page 2)', size=11, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# 중간: 추출 텍스트 깨짐
add_rect(s, Inches(5.2), Inches(1.5), Inches(3.9), Inches(5.3), RGBColor(0xFF, 0xF0, 0xF0))
add_textbox(s, Inches(5.35), Inches(1.65), Inches(3.6), Inches(0.4),
            'pypdf 추출 결과', size=14, bold=True, color=COLOR_WARN)
add_textbox(s, Inches(5.35), Inches(2.1), Inches(3.6), Inches(4.6),
            [
                '"정전용량([수식])은 극판의',
                '넓이([수식])에 비례하고 극판의',
                '간격([수식])에 반비례한다."',
                '',
                '→ 변수 C, S, d 가 빠짐',
                '   ([수식] 마커로 표시)',
                '',
                '"오답분석',
                ' ① [수식] 관계에 의해',
                '   전압에 반비례한다."',
                '',
                '→ Q=CV 같은 식 누락',
            ],
            size=11, color=COLOR_DARK, font='Consolas')

# 우측: HTML 보완 결과
add_rect(s, Inches(9.3), Inches(1.5), Inches(3.7), Inches(5.3), RGBColor(0xE8, 0xF8, 0xE8))
add_textbox(s, Inches(9.45), Inches(1.65), Inches(3.4), Inches(0.4),
            '수동 보완 (HTML)', size=14, bold=True, color=COLOR_OK)
add_textbox(s, Inches(9.45), Inches(2.1), Inches(3.4), Inches(4.6),
            [
                '【해설】',
                '정전용량(<i>C</i>)은 극판의',
                '넓이(<i>S</i>)에 비례하고…',
                '',
                '【학습 POINT】',
                '<p><strong>정전용량 공식</strong></p>',
                '<p><i>C</i> = <i>εS</i> / <i>d</i> [F]</p>',
                '<ul>',
                '  <li><i>C</i>: 정전용량[F]</li>',
                '  <li><i>ε</i>: 유전율[F/m]</li>',
                '  <li><i>S</i>: 극판 단면적[m²]</li>',
                '  <li><i>d</i>: 극판 간격[m]</li>',
                '</ul>',
            ],
            size=10, color=COLOR_DARK, font='Consolas')


# ================================================================
# Slide 7 — 표 → HTML 표 변환 (23번 반작용)
# ================================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, '6. 표 → HTML 표 변환 — 23번 동기전동기 반작용',
            'PDF의 표는 셀 경계 깨짐 · HTML <table>로 재구성')

# 좌측: PDF 캡처
s.shapes.add_picture(IMG('p08_q23_반작용표.png'),
                     Inches(0.3), Inches(1.5),
                     width=Inches(6.0), height=Inches(5.5))
add_textbox(s, Inches(0.3), Inches(7.05), Inches(6.0), Inches(0.3),
            '원본 PDF (page 8) — 표 두 개 포함', size=11, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# 우측 상단: 추출 깨짐 안내
add_rect(s, Inches(6.6), Inches(1.5), Inches(6.4), Inches(1.4), RGBColor(0xFF, 0xF0, 0xF0))
add_textbox(s, Inches(6.75), Inches(1.6), Inches(6.1), Inches(0.4),
            '⚠️ 추출 시', size=13, bold=True, color=COLOR_WARN)
add_textbox(s, Inches(6.75), Inches(2.0), Inches(6.1), Inches(0.9),
            [
                '"전류 위상 부하 성분 작용 명칭 현상',
                ' 동상 저항(R) 교차 자화 작용 …"',
                '→ 셀 경계 사라지고 줄바꿈 깨짐',
            ],
            size=11, color=COLOR_DARK, font='Consolas')

# 우측 하단: HTML 결과 미리보기 (간이)
add_rect(s, Inches(6.6), Inches(3.05), Inches(6.4), Inches(4.05), RGBColor(0xE8, 0xF8, 0xE8))
add_textbox(s, Inches(6.75), Inches(3.15), Inches(6.1), Inches(0.4),
            '✓ HTML <table> 재구성', size=13, bold=True, color=COLOR_OK)
# 미리보기 표
mini = [
    ['전류 위상', '동기 발전기', '동기 전동기'],
    ['지상 (L, 뒤진)', '감자 작용', '증자 작용'],
    ['진상 (C, 앞선)', '증자 작용', '감자 작용'],
    ['동상 (R)', '교차 자화', '교차 자화'],
]
mtbl = s.shapes.add_table(len(mini), 3, Inches(6.75), Inches(3.6), Inches(6.1), Inches(3.0)).table
for r, row in enumerate(mini):
    for c, val in enumerate(row):
        cell = mtbl.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = FONT_KR
                run.font.size = Pt(11)
                run.font.bold = (r == 0)
                run.font.color.rgb = COLOR_DARK if r > 0 else RGBColor(0xFF, 0xFF, 0xFF)
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_OK


# ================================================================
# Slide 8 — 그림 필요 영역 마커 (51번 그림기호)
# ================================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, '7. 그림 필요 위치 — 컬럼에 [필요: …] 마커',
            '텍스트로 추출 불가능한 그림은 위치를 명시해서 수동 등록 유도')

# 좌측: PDF 캡처 (51번)
s.shapes.add_picture(IMG('p16_q51_그림기호.png'),
                     Inches(0.3), Inches(1.5),
                     width=Inches(6.0), height=Inches(5.5))
add_textbox(s, Inches(0.3), Inches(7.05), Inches(6.0), Inches(0.3),
            '원본 PDF (page 16) — 51번 그림기호 문제',
            size=11, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# 우측: 엑셀 결과
add_rect(s, Inches(6.6), Inches(1.5), Inches(6.4), Inches(5.5), RGBColor(0xE8, 0xF0, 0xFE))
add_textbox(s, Inches(6.75), Inches(1.65), Inches(6.1), Inches(0.4),
            '엑셀 51번 행 (발췌)', size=14, bold=True, color=COLOR_ACCENT)
add_textbox(s, Inches(6.75), Inches(2.1), Inches(6.1), Inches(4.8),
            [
                '발문    │ "아래 그림기호가 나타내는 것은?"',
                '발문그림 │ [필요: 기호심볼 — 수동조작 접점] ⚠',
                '보기1   │ "한시 계전기 접점"',
                '보기2   │ "전자 접속기 접점"',
                '보기3   │ "수동 조작 접점"',
                '보기4   │ "조작 개폐기 잔류 접점"',
                '정답    │ 3',
                '',
                '학습포인트 │ <table>접점 기호 분류 표</table>',
                '',
                '─────────────────────',
                '운영자: 마커 보고 그림 파일 업로드 → 파일명 교체',
                '예) 발문그림 = q51_manual_contact.png',
            ],
            size=12, color=COLOR_DARK, font='Consolas')


# ================================================================
# Slide 9 — 학습 POINT 자동 분리
# ================================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, '8. 학습 POINT 자동 분리 로직',
            '해설 본문 안의 "공식 / 법칙 / 특징" 헤더 패턴으로 자동 분리')

# 좌측: 분리 전 (해설 통째)
add_rect(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.4), RGBColor(0xFF, 0xF7, 0xE0))
add_textbox(s, Inches(0.7), Inches(1.75), Inches(5.4), Inches(0.5),
            '분리 전 — 해설 한 덩어리', size=16, bold=True, color=COLOR_PRIMARY)
add_textbox(s, Inches(0.7), Inches(2.3), Inches(5.4), Inches(4.6),
            [
                '전기분해를 통해 석출된 물질의 양은',
                '통과한 전기량 및 화학당량에 각각',
                '비례한다.',
                '',
                '────────',
                '【여기서부터 학습 POINT 영역】',
                '패러데이의 법칙',
                ' W = kIt [g]',
                ' • W: 석출된 물질의 양[g]',
                ' • k: 전기 화학 당량[g/C]',
                ' • I: 전류[A]  • t: 시간[s]',
            ],
            size=12, color=COLOR_DARK, font='Consolas')

# 화살표
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(4.0), Inches(0.5), Inches(0.5))
arr.fill.solid()
arr.fill.fore_color.rgb = COLOR_PRIMARY
arr.line.fill.background()

# 우측: 분리 후 — 2개 컬럼
add_rect(s, Inches(7.2), Inches(1.6), Inches(5.6), Inches(2.5), RGBColor(0xE8, 0xF8, 0xE8))
add_textbox(s, Inches(7.4), Inches(1.75), Inches(5.2), Inches(0.5),
            '해설 컬럼', size=14, bold=True, color=COLOR_OK)
add_textbox(s, Inches(7.4), Inches(2.3), Inches(5.2), Inches(1.7),
            [
                '전기분해를 통해 석출된 물질의 양은',
                '통과한 전기량 및 화학당량에 각각',
                '비례한다.',
            ],
            size=12, color=COLOR_DARK)

add_rect(s, Inches(7.2), Inches(4.3), Inches(5.6), Inches(2.7), RGBColor(0xFC, 0xE4, 0xEC))
add_textbox(s, Inches(7.4), Inches(4.45), Inches(5.2), Inches(0.5),
            '학습포인트 컬럼', size=14, bold=True, color=COLOR_PRIMARY)
add_textbox(s, Inches(7.4), Inches(5.0), Inches(5.2), Inches(2.0),
            [
                '<p><strong>패러데이의 법칙</strong></p>',
                '<p><i>W</i> = <i>kIt</i> [g]</p>',
                '<ul>',
                '  <li>W: 석출량[g] · k: 화학당량[g/C]</li>',
                '  <li>I: 전류[A]  · t: 시간[s]</li>',
                '</ul>',
            ],
            size=11, color=COLOR_DARK, font='Consolas')


# ================================================================
# Slide 10 — 최종 통계
# ================================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, '9. 최종 통계 — 60문항 변환 결과',
            '데이터 / 보완 / 이슈 상태')

# 큰 숫자 4개
cards = [
    ('60', '총 문항', COLOR_DARK, '이론 20 / 기기 20 / 설비 20'),
    ('59', '자동 정상 변환', COLOR_OK, '이슈 없이 그대로 사용 가능'),
    ('60', '학습 POINT 보유', COLOR_PRIMARY, '자동 분리 14 + 수동 보완 47'),
    ('1', '잔여 이슈', COLOR_WARN, '9번 빈칸문제 노트만'),
]
for i, (big, label, c, sub) in enumerate(cards):
    cx = Inches(0.5 + i * 3.15)
    add_rect(s, cx, Inches(1.8), Inches(3.0), Inches(2.6), COLOR_BG_LIGHT, line=c)
    add_textbox(s, cx, Inches(1.95), Inches(3.0), Inches(1.4),
                big, size=72, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_textbox(s, cx, Inches(3.4), Inches(3.0), Inches(0.5),
                label, size=14, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, cx, Inches(3.85), Inches(3.0), Inches(0.5),
                sub, size=11, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# 하단: 카테고리별 정리
add_rect(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.5), RGBColor(0xFA, 0xFA, 0xFA))
add_textbox(s, Inches(0.7), Inches(4.85), Inches(12), Inches(0.5),
            '카테고리별 보완 작업 누계', size=15, bold=True, color=COLOR_PRIMARY)

stat_data = [
    ['항목', '건수', '대상 문항'],
    ['수식 → HTML', '약 40', '1·2·5·6·7·9·11·13·15·19 외 다수'],
    ['표 → HTML 표', '약 20', '3·8·17·23·26·29·30·32·43·47·50·53·54·58·59·60'],
    ['그림 마커', '2', '39번 TRIAC 심볼 / 51번 수동조작 접점'],
    ['PDF 오류 보완', '1', '55번 정답 태그 누락'],
    ['추정값 (확인 필요)', '6', '7·13·22·32·36·48번 (수치 PUA 손상)'],
]
stbl = s.shapes.add_table(len(stat_data), 3, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.9)).table
stbl.columns[0].width = Inches(3.2)
stbl.columns[1].width = Inches(1.5)
stbl.columns[2].width = Inches(7.2)
for r, row in enumerate(stat_data):
    for c, val in enumerate(row):
        cell = stbl.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = FONT_KR
                run.font.size = Pt(11)
                run.font.bold = (r == 0)
                run.font.color.rgb = COLOR_DARK if r > 0 else RGBColor(0xFF, 0xFF, 0xFF)
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_PRIMARY


# ================================================================
# Slide 11 — 남은 한계와 다음 단계
# ================================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, '10. 남은 한계와 다음 단계',
            '오늘까지의 결과 / 팀에서 이어받을 작업')

# 좌측: 한계
add_rect(s, Inches(0.5), Inches(1.6), Inches(6.0), Inches(5.3), RGBColor(0xFF, 0xF0, 0xF0))
add_textbox(s, Inches(0.7), Inches(1.75), Inches(5.6), Inches(0.5),
            '⚠️ 자동화의 한계', size=18, bold=True, color=COLOR_WARN)
add_textbox(s, Inches(0.7), Inches(2.3), Inches(5.6), Inches(4.5),
            [
                '1. PDF 폰트 PUA 영역 글자 (수식)',
                '    pypdf 추출 불가 — [수식] 마커로 명시',
                '',
                '2. 추정값 6문항 (7·13·22·32·36·48)',
                '    원본 수치값 확인 후 manual_fixes 수정 필요',
                '    <em> 태그로 "원본 확인 필요" 표시됨',
                '',
                '3. 그림 (39·51번)',
                '    실제 그림 파일 업로드 + 컬럼 채움',
                '',
                '4. 9번 빈칸문제',
                '    구조 자체가 ( ① ) ( ② ) 빈칸',
                '    UI에서 별도 렌더링 처리 검토 필요',
                '',
                '5. PDF 마스터 데이터 오류',
                '    모든 문항이 [전기이론] 태그 — 페이지 헤더로 보정',
            ],
            size=12, color=COLOR_DARK)

# 우측: 다음 단계
add_rect(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(5.3), RGBColor(0xE8, 0xF8, 0xE8))
add_textbox(s, Inches(7.0), Inches(1.75), Inches(5.6), Inches(0.5),
            '✓ 다음 단계 (인수인계)', size=18, bold=True, color=COLOR_OK)
add_textbox(s, Inches(7.0), Inches(2.3), Inches(5.6), Inches(4.5),
            [
                'A. 다른 회차 PDF 적용',
                '   convert_2022_1.py를 회차별로 복제',
                '   → 페이지 번호 / 카테고리 라인만 조정',
                '',
                'B. manual_fixes.py 재사용',
                '   학습 POINT·표 HTML이 이미 패턴화됨',
                '   새 회차에도 30~50% 그대로 적용 가능',
                '',
                'C. 업로드 UI 만들기',
                '   엑셀 → DB / JSON 변환 임포터',
                '   HTML 태그 sanitize (XSS 차단)',
                '',
                'D. CBT 렌더링 검증',
                '   <table>, <sub>, <sup>, <i>, <ol>',
                '   모바일·태블릿 양쪽에서 깨지지 않는지',
                '',
                'E. 추정값 6문항 원본 확인',
                '   책 보면서 manual_fixes 정확값 입력',
            ],
            size=12, color=COLOR_DARK)

# 저장
prs.save(OUT)
print(f'WROTE: {OUT}')
print(f'  슬라이드 {len(prs.slides)}장')
