"""엑셀 v3 ↔ 등록 페이지 연동 설명 PPT v2 — 슬라이드 3 확장.
변경:
  - 슬라이드 3 (전체 엑셀) — 각 영역 한 줄 설명 추가
  - 슬라이드 4 신규 — 메타 영역 클로즈업 (좌반)
  - 슬라이드 5 신규 — 콘텐츠 영역 클로즈업 (우반, 블록 시퀀스 강조)
  - 이전 "메타 4영역 카드" 슬라이드는 4·5에 흡수
출력: data/presentation/엑셀-등록페이지_연동.pptx (덮어쓰기)
"""
import os, sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = lambda f: os.path.join('data/presentation/img', f)
OUT = 'data/presentation/엑셀-등록페이지_연동.pptx'
os.makedirs(os.path.dirname(OUT), exist_ok=True)

FONT = '맑은 고딕'
MONO = 'Consolas'

C_PRIMARY = RGBColor(0xC2, 0x18, 0x5B)
C_ACCENT = RGBColor(0x21, 0x96, 0xF3)
C_DARK = RGBColor(0x1F, 0x29, 0x37)
C_MUTED = RGBColor(0x6B, 0x72, 0x80)
C_BG = RGBColor(0xFA, 0xFA, 0xFA)
C_OK = RGBColor(0x05, 0x96, 0x69)
C_WARN = RGBColor(0xD9, 0x77, 0x06)

# 영역별 (엑셀과 동일)
C_META = RGBColor(0x47, 0x55, 0x69)
C_CODE = RGBColor(0x4F, 0x46, 0xE5)
C_CLASS = RGBColor(0x08, 0x91, 0xB2)
C_ATTR = RGBColor(0xE1, 0x1D, 0x48)
C_STEM = RGBColor(0xDB, 0x27, 0x77)
C_OPT = RGBColor(0x25, 0x63, 0xEB)
C_IMG = RGBColor(0xD9, 0x77, 0x06)
C_ANS = RGBColor(0x0D, 0x94, 0x88)
C_EXP = RGBColor(0xEA, 0x58, 0x0C)
C_WR = RGBColor(0x93, 0x33, 0xEA)
C_LP = RGBColor(0x05, 0x96, 0x69)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def tb(slide, l, t, w, h, lines, *, size=14, bold=False, color=C_DARK,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return box


def rect(slide, l, t, w, h, fill, line=None, line_w=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def title(slide, t, sub=None):
    rect(slide, 0, 0, prs.slide_width, Inches(0.15), C_PRIMARY)
    tb(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
       t, size=26, bold=True)
    if sub:
        tb(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4),
           sub, size=13, color=C_MUTED)


def chip(slide, l, t, color, label, w=Inches(2.5), label_size=12):
    rect(slide, l, t, Inches(0.2), Inches(0.3), color)
    tb(slide, l + Inches(0.3), t - Inches(0.02), w, Inches(0.4),
       label, size=label_size, bold=True)


def code_box(slide, l, t, w, h, lines):
    rect(slide, l, t, w, h, RGBColor(0xF5, 0xF5, 0xF5))
    tb(slide, l + Inches(0.1), t + Inches(0.08), w - Inches(0.2), h - Inches(0.15),
       lines, size=11, font=MONO)


def region_card(slide, l, t, w, h, color, emoji, name, body_lines,
                title_size=14, body_size=11):
    """영역 설명 카드 — 색 띠 + 제목 + 본문."""
    rect(slide, l, t, w, h, C_BG, line=color, line_w=2)
    rect(slide, l, t, w, Inches(0.42), color)
    tb(slide, l + Inches(0.15), t + Inches(0.06), w - Inches(0.3), Inches(0.35),
       f'{emoji}  {name}', size=title_size, bold=True,
       color=RGBColor(0xFF, 0xFF, 0xFF))
    tb(slide, l + Inches(0.15), t + Inches(0.5), w - Inches(0.3), h - Inches(0.55),
       body_lines, size=body_size, color=C_DARK)


# =============================================================
# Slide 1 — 표지
# =============================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, prs.slide_width, prs.slide_height, C_BG)
rect(s, Inches(0.6), Inches(2.8), Inches(0.5), Inches(0.06), C_PRIMARY)
tb(s, Inches(0.6), Inches(2.95), Inches(12), Inches(1.2),
   '엑셀 v3 양식 ↔ 등록 페이지 연동', size=42, bold=True)
tb(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.6),
   '같은 데이터 모델 / 양방향 변환 / HTML-native 직렬화',
   size=18, color=C_MUTED)
tb(s, Inches(0.6), Inches(5.0), Inches(12), Inches(1.5),
   [
       '· 엑셀 한 행 = 등록 페이지 한 화면 = 시험 한 문항',
       '· 영역별 컬러 그룹화 — 출처 / 코드 / 분류 / 속성 / 콘텐츠',
       '· 해설·오답분석·학습 POINT는 HTML 작성 · <img>로 그림 자유 배치',
       '· xlsx ↔ JSON 양방향 변환 검증 완료 (왕복 일치 ✅)',
   ], size=14)


# =============================================================
# Slide 2 — 전체 워크플로
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '1. 전체 워크플로 — 엑셀 ↔ 등록 페이지 ↔ 시험 페이지',
      '같은 데이터 모델을 세 곳에서 공유')

boxes = [
    ('📊  엑셀 v3 양식', C_META, [
        'data/templates/', 'question-template-v3.xlsx', '',
        '- 33컬럼, 그룹 헤더', '- 영역별 컬러',
        '- 드롭다운 검증', '- 배치 등록용',
    ]),
    ('🖊  등록 페이지', C_STEM, [
        'src/app/admin/register', '/admin/register', '',
        '- 실시간 미리보기', '- HTML 입력',
        '- 블록 시퀀스 UI', '- 단일 문항 등록용',
    ]),
    ('📝  시험 페이지', C_ACCENT, [
        '(향후 구현)', '문항코드 순 정렬', '',
        '- 회차/연도 필터', '- 빈출도/난이도 필터',
        '- 동일 데이터 모델', '- 학습자에게 표시',
    ]),
]
yy = Inches(1.7)
w = Inches(4.0)
for i, (t, color, lines) in enumerate(boxes):
    cx = Inches(0.5 + i * (4.0 + 0.15))
    rect(s, cx, yy, w, Inches(4.0), C_BG, line=color, line_w=3)
    rect(s, cx, yy, w, Inches(0.6), color)
    tb(s, cx, yy + Inches(0.1), w, Inches(0.4),
       t, size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
       align=PP_ALIGN.CENTER)
    tb(s, cx + Inches(0.2), yy + Inches(0.8), w - Inches(0.4), Inches(3.0),
       lines, size=12, color=C_DARK, font=MONO)

for i in range(2):
    cx = Inches(0.5 + (i + 1) * 4.0 + i * 0.15 + 0.01)
    arr = s.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW,
                              cx, yy + Inches(1.7), Inches(0.13), Inches(0.5))
    arr.fill.solid(); arr.fill.fore_color.rgb = C_MUTED
    arr.line.fill.background()

rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2),
     RGBColor(0xFF, 0xF7, 0xE0))
tb(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.5),
   '🔄 양방향 변환: tools/question_io.py', size=14, bold=True, color=C_WARN)
tb(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.7),
   [
       'python tools/question_io.py xlsx2json input.xlsx output.json',
       'python tools/question_io.py json2xlsx input.json output.xlsx',
   ], size=12, font=MONO)


# =============================================================
# Slide 3 — 엑셀 v3 양식 전체 (영역별 한 줄 설명 추가)
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '2. 엑셀 v3 양식 — 33컬럼 한눈에 (영역별 그룹 헤더)',
      'A3 가로 한 페이지 · 1행 그룹 헤더 + 2행 컬럼 헤더 + 데이터 행')

s.shapes.add_picture(IMG('xlsx_question-template-v3-new_p1.png'),
                     Inches(0.3), Inches(1.4), width=Inches(12.7))

# 하단 영역 한 줄 설명 (2단)
brief = [
    (C_META, '출처', '시험 출신 (회차)'),
    (C_CODE, '코드', '정렬·식별 키'),
    (C_CLASS, '분류', '카테고리 (세로축)'),
    (C_ATTR, '속성', '학습 메타 (가로축)'),
    (C_STEM, '발문', '문제 본문'),
    (C_OPT, '보기', '4지선다'),
    (C_ANS, '정답', '1·2·3·4'),
    (C_EXP, '해설', 'HTML + <img>'),
    (C_WR, '오답분석', 'HTML <ol> + <img>'),
    (C_LP, '학습POINT', 'HTML + <img>'),
]
# 2줄 5컬럼 그리드
y0 = Inches(6.55)
col_w = Inches(2.55)
for i, (color, key, desc) in enumerate(brief):
    col = i % 5
    row = i // 5
    cx = Inches(0.4 + col * 2.55)
    cy = y0 + row * Inches(0.4)
    rect(s, cx, cy, Inches(0.15), Inches(0.32), color)
    tb(s, cx + Inches(0.22), cy - Inches(0.02), col_w - Inches(0.3), Inches(0.4),
       [f'{key}  ', f'{desc}'], size=10, color=C_DARK)
    tb(s, cx + Inches(0.22), cy - Inches(0.02), Inches(1.0), Inches(0.18),
       key, size=11, bold=True, color=color)
    tb(s, cx + Inches(0.22), cy + Inches(0.16), col_w - Inches(0.3), Inches(0.18),
       desc, size=9, color=C_MUTED)


# =============================================================
# Slide 4 — 메타 영역 클로즈업 (좌반 17컬럼) + 4개 영역 카드
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '3. 메타 영역 — 출처 / 코드 / 분류(세로축) / 속성(가로축)',
      '엑셀 좌반 17컬럼 · 시험 출신·정렬·카테고리·학습 메타')

# 상단: 좌반 클로즈업 이미지
s.shapes.add_picture(IMG('xlsx_left_meta.png'),
                     Inches(0.3), Inches(1.4), width=Inches(12.7))

# 하단: 4개 영역 카드 (2x2)
card_y = Inches(4.7)
card_w = Inches(6.2)
card_h = Inches(1.35)

region_card(s, Inches(0.4), card_y, card_w, card_h, C_META,
            '📌', '출처 (6컬럼) — 과정 · 연도 · 회차 · 번호 · 사용교재 · 교재구분',
            [
                '· 어떤 시험에서 나온 문제인지 식별',
                '· 같은 회차 문항 묶기 / 연도·회차별 필터링 기준',
                '· 사용교재·교재구분: 변형이력 추적 시 기준 (다른 판 비교)',
            ])

region_card(s, Inches(6.8), card_y, card_w, card_h, C_CODE,
            '🔖', '코드 (1컬럼) — 문항코드',
            [
                '· 정렬 키 — 시험 페이지가 이 코드 순으로 문항 나열',
                '· 규칙: elec_<교재구분>_<연도>_<회차>_<일련번호>',
                '· 회차 내 고유 · 다른 회차와 충돌 없음',
            ])

region_card(s, Inches(0.4), card_y + Inches(1.45), card_w, card_h, C_CLASS,
            '🗂', '분류 (5컬럼, 세로축) — 과목ID · 챕터 · 대유형 · 중유형 · 내용',
            [
                '· 콘텐츠 카테고리. 5단계 트리 (좁아질수록 구체적)',
                '· 학습 시 챕터/유형별 검색·필터 키',
                '· 통계: 어느 챕터·유형이 약점인지 산출 기준',
            ])

region_card(s, Inches(6.8), card_y + Inches(1.45), card_w, card_h, C_ATTR,
            '🏷', '속성 (5컬럼, 가로축) — 빈출도 · 난이도 · 문제유형 · 변형이력 · 비고',
            [
                '· 빈출도·난이도 ★1~5 — 시험 페이지 필터·추천 기준',
                '· 문제유형 (단답형/계산형/도식형/조합형/빈칸형/완성형)',
                '· 변형이력 0=원본 · 1·2…=변형문제 · 비고는 자유 메모',
            ])


# =============================================================
# Slide 5 — 콘텐츠 영역 클로즈업 (우반 15컬럼) + 6개 영역 카드
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '4. 콘텐츠 영역 — 발문·보기·정답·해설·오답분석·학습POINT',
      '엑셀 우반 15컬럼 · 해설·오답분석·학습POINT는 한 셀에 HTML 작성 (<img>로 그림)')

# 상단: 우반 클로즈업 이미지 (가로로 김)
s.shapes.add_picture(IMG('xlsx_right_content.png'),
                     Inches(0.2), Inches(1.4), width=Inches(12.9))

# 하단: 6개 영역 카드 (3x2)
card_y = Inches(4.4)
card_w = Inches(4.2)
card_h = Inches(1.4)
gap = Inches(0.15)

region_card(s, Inches(0.4), card_y, card_w, card_h, C_STEM,
            '📝', '발문 (3컬럼)',
            [
                '· 발문 / 조건 / 발문그림',
                '· HTML + LaTeX 수식 ($…$ 인라인 / $$…$$ 디스플레이)',
                '· 그림 필요 시 [필요: 회로도] 마커',
            ], title_size=12)

region_card(s, Inches(4.75), card_y, card_w, card_h, C_OPT,
            '✏', '보기 (8컬럼)',
            [
                '· 보기1~4 + 보기1~4그림 각각',
                '· 텍스트형 / 그림형(회로도 4개 비교 등) 둘 다',
                '· 보기 자체가 그림이면 보기N그림 컬럼 사용',
            ], title_size=12)

region_card(s, Inches(9.1), card_y, card_w, card_h, C_ANS,
            '✓', '정답 (1컬럼)',
            [
                '· 1·2·3·4 중 하나의 정수',
                '· 시험 페이지 채점 키',
                '· 드롭다운으로 입력 (오타 방지)',
            ], title_size=12)

region_card(s, Inches(0.4), card_y + Inches(1.5), card_w, card_h, C_EXP,
            '💡', '해설 (1컬럼) — HTML',
            [
                '· 한 셀에 HTML 작성 (텍스트·그림 교차)',
                '· 그림: <img src="파일명"> 본문 삽입',
                '· 그림 여러 개·위치 제한 없음',
                '· 글-그림-글-그림 패턴 가능 (PDF 순서 보존)',
            ], title_size=12, body_size=10)

region_card(s, Inches(4.75), card_y + Inches(1.5), card_w, card_h, C_WR,
            '🔍', '오답분석 (1컬럼)',
            [
                '· HTML <ol><li>...</li></ol> 권장',
                '· 정답이 ②면 <li value="2">부터 시작 가능',
                '· 그림 필요 시 <img src="파일명"> 삽입 가능',
            ], title_size=12)

region_card(s, Inches(9.1), card_y + Inches(1.5), card_w, card_h, C_LP,
            '🎯', '학습 POINT (1컬럼) — HTML',
            [
                '· 해설과 동일 (HTML + <img>)',
                '· 공식·법칙·요약 표 (<table>) 위주',
                '· 표는 <table>로 작성',
                '· 그림 미정 시 <img src="[필요: 설명]">',
            ], title_size=12, body_size=10)


# =============================================================
# Slide 6 — 등록 페이지 빈 상태
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '5. 등록 페이지 — 빈 입력 화면',
      'URL: /admin/register · 좌측 입력 / 우측 실시간 미리보기')
s.shapes.add_picture(IMG('register_empty.png'),
                     Inches(0.5), Inches(1.4), width=Inches(12.3))


# =============================================================
# Slide 7 — 01번 매칭 (등록 페이지 미리보기 ↔ 엑셀 행 표)
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '6. 매칭 — 01번 콘덴서 (등록 페이지 ↔ 엑셀 한 행)',
      '같은 데이터가 두 형태로 보임')

tb(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(0.4),
   '🖊 등록 페이지 (미리보기 영역)', size=14, bold=True, color=C_STEM)
s.shapes.add_picture(IMG('register_q01_preview.png'),
                     Inches(0.4), Inches(1.9), width=Inches(6.0))

tb(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(0.4),
   '📊 엑셀 v3 — 같은 데이터의 엑셀 행', size=14, bold=True, color=C_META)
fields = [
    ('영역', '컬럼', '값'),
    ('출처', '연도 / 회차', '2022 / 1'),
    ('코드', '문항코드', 'elec_A_2022_01_01'),
    ('분류', '과목ID', 'theory'),
    ('분류', '챕터 → 내용', '정전기 › 콘덴서 › 정전용량'),
    ('속성', '빈출도 / 난이도', '4 / 2'),
    ('속성', '문제유형', '단답형'),
    ('발문', '발문', '콘덴서의 정전용량에 대한 설명으로 틀린 것은?'),
    ('보기', '보기1~4', '①전압 ②전하 ③넓이 ④간격'),
    ('정답', '정답(1~4)', '4'),
    ('해설', '해설', '정전용량(C)은 극판의...'),
    ('학습POINT', '학습포인트', '<table>정전용량 공식</table>'),
]
tbl = s.shapes.add_table(len(fields), 3, Inches(6.9), Inches(1.9),
                          Inches(6.0), Inches(4.5)).table
tbl.columns[0].width = Inches(1.1)
tbl.columns[1].width = Inches(1.6)
tbl.columns[2].width = Inches(3.3)
for r, row in enumerate(fields):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = str(val)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = FONT
                run.font.size = Pt(10)
                run.font.bold = (r == 0)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r == 0 else C_DARK
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_META


# =============================================================
# Slide 8 — 블록 시퀀스 입력 UI
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '7. 블록 시퀀스 입력 — 해설·학습 POINT (등록 페이지)',
      '텍스트 / 그림 블록을 + 버튼으로 자유 배치 · 순서 변경 / 삭제')

s.shapes.add_picture(IMG('register_q21_explain_blocks.png'),
                     Inches(0.4), Inches(1.5), width=Inches(7.0))
tb(s, Inches(0.4), Inches(7.05), Inches(7.0), Inches(0.3),
   '21번 차동계전기 — 해설 영역 (텍스트 → 그림 → 텍스트 3블록)',
   size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

right = Inches(7.6)
rect(s, right, Inches(1.5), Inches(5.4), Inches(5.7), C_BG, line=C_EXP)
tb(s, right + Inches(0.2), Inches(1.65), Inches(5.0), Inches(0.4),
   '🧩 블록 시퀀스 UI', size=15, bold=True, color=C_EXP)
tb(s, right + Inches(0.2), Inches(2.1), Inches(5.0), Inches(5.0),
   [
       '◆ 블록 종류',
       '   T 텍스트 — HTML 허용',
       '   🖼 그림 — 파일명 또는 [필요: …]',
       '',
       '◆ 추가',
       '   • 영역 하단의 [+ 텍스트] / [+ 그림]',
       '   • 블록 사이 [+ 아래에 텍스트 / 그림]',
       '',
       '◆ 조작',
       '   • ↑ ↓ 위/아래 이동',
       '   • ✕ 삭제',
       '',
       '◆ 동일 패턴 적용',
       '   • 해설',
       '   • 학습 POINT',
       '',
       '◆ 활용 예',
       '   본문 → 회로도 → 추가 설명 → 표',
       '   (PDF 페이지 시퀀스 그대로 옮기기 가능)',
   ], size=12)


# =============================================================
# Slide 9 — HTML-native 직렬화 (JSON ↔ 엑셀 셀)
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '8. 엑셀 셀 직렬화 — HTML-native',
      '<img src="파일명">로 그림 삽입 · 파서가 <img> 경계로 블록 자동 분리')

tb(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(0.4),
   '🖊 등록 페이지 JSON 형태', size=14, bold=True, color=C_STEM)
code_box(s, Inches(0.4), Inches(1.85), Inches(6.2), Inches(5.0),
         [
             '"해설블록": [',
             '  {',
             '    "type": "text",',
             '    "value": "변압기 내부 권선의',
             '       층간 단락이나 지락 사고 시..."',
             '  },',
             '  {',
             '    "type": "image",',
             '    "value": "q21_differential_diagram.png"',
             '  },',
             '  {',
             '    "type": "text",',
             '    "value": "이 차이를 검출하여 동작하는',
             '       <strong>차동 계전기</strong>가',
             '       가장 적합하다."',
             '  }',
             ']',
         ])

arr = s.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW,
                         Inches(6.7), Inches(4.0),
                         Inches(0.2), Inches(0.6))
arr.fill.solid(); arr.fill.fore_color.rgb = C_OK
arr.line.fill.background()
tb(s, Inches(6.5), Inches(4.65), Inches(0.7), Inches(0.3),
   '변환', size=10, color=C_OK, align=PP_ALIGN.CENTER)

tb(s, Inches(7.1), Inches(1.4), Inches(6.0), Inches(0.4),
   '📊 엑셀 "해설" 셀 하나에 직렬화', size=14, bold=True, color=C_META)
code_box(s, Inches(7.1), Inches(1.85), Inches(6.0), Inches(5.0),
         [
             '변압기 내부 권선의 층간 단락이나',
             '지락 사고 시, 입력 전류와 출력',
             '전류 사이에 차이가 발생한다.',
             '',
             '<img src="q21_differential_diagram.png">',
             '',
             '이 차이를 검출하여 동작하는',
             '<strong>차동 계전기</strong>가',
             '가장 적합하다.',
         ])

rect(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
     RGBColor(0xFF, 0xF7, 0xE0))
tb(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
   '규칙: 1) HTML 그대로 작성  2) 그림은 <img src="파일명">  3) 개수·위치 제한 없음 (구 [IMG:]도 호환)',
   size=11, color=C_WARN, bold=True, align=PP_ALIGN.CENTER)


# =============================================================
# Slide 10 — 21번 풀 화면
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '9. 실제 입력 화면 — 21번 차동계전기 (블록 3개)',
      '좌측 입력 / 우측 실시간 미리보기 — 즉시 검증 가능')
s.shapes.add_picture(IMG('register_q21_blocks.png'),
                     Inches(0.3), Inches(1.4), width=Inches(12.7))


# =============================================================
# Slide 11 — 양방향 변환 CLI
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '10. 양방향 변환 라이브러리 — tools/question_io.py',
      '엑셀 ↔ JSON 자동 변환 · 왕복 일치 검증 완료')

rect(s, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.5), C_BG, line=C_ACCENT)
tb(s, Inches(0.6), Inches(1.65), Inches(5.8), Inches(0.5),
   '📥 엑셀 → JSON (xlsx2json)', size=15, bold=True, color=C_ACCENT)
tb(s, Inches(0.6), Inches(2.15), Inches(5.8), Inches(0.4),
   'PDF 변환 결과나 배치 등록 엑셀을 등록 페이지 형식으로',
   size=11, color=C_MUTED)
code_box(s, Inches(0.6), Inches(2.6), Inches(5.8), Inches(1.2),
         [
             '$ python tools/question_io.py \\',
             '    xlsx2json input.xlsx \\',
             '    output.json',
         ])
tb(s, Inches(0.6), Inches(3.95), Inches(5.8), Inches(0.4),
   '처리 흐름:', size=12, bold=True)
tb(s, Inches(0.6), Inches(4.35), Inches(5.8), Inches(2.5),
   [
       '1. 문항등록 시트 헤더 읽기',
       '2. 각 행을 dict로 변환',
       '3. 해설/학습포인트 컬럼 → 블록 배열로 파싱',
       '4. <img> / 레거시 [IMG:] 경계 분리',
       '5. 블록에 UUID id 부여',
       '6. JSON 배열로 저장',
   ], size=11)

rect(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.5), C_BG, line=C_OK)
tb(s, Inches(7.0), Inches(1.65), Inches(5.8), Inches(0.5),
   '📤 JSON → 엑셀 (json2xlsx)', size=15, bold=True, color=C_OK)
tb(s, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.4),
   '등록 페이지 JSON을 배치 엑셀로',
   size=11, color=C_MUTED)
code_box(s, Inches(7.0), Inches(2.6), Inches(5.8), Inches(1.2),
         [
             '$ python tools/question_io.py \\',
             '    json2xlsx items.json \\',
             '    output.xlsx',
         ])
tb(s, Inches(7.0), Inches(3.95), Inches(5.8), Inches(0.4),
   '처리 흐름:', size=12, bold=True)
tb(s, Inches(7.0), Inches(4.35), Inches(5.8), Inches(2.5),
   [
       '1. JSON 배열 로드',
       '2. 각 항목을 v3 양식 컬럼으로 매핑',
       '3. 해설블록/학습포인트블록 → 한 셀에 직렬화',
       '4. 그림 블록 → <img src="파일명"> 자동 삽입',
       '5. 텍스트·그림 블록 사이 빈 줄 구분',
       '6. xlsx 저장',
   ], size=11)


# =============================================================
# Slide 12 — 다음 단계 체크리스트
# =============================================================
s = prs.slides.add_slide(BLANK)
title(s, '11. 다음 단계 — 시험 페이지 / 임포터 / 일괄 등록',
      '데이터 모델은 확정 — 이걸로 시스템 연결')

rect(s, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.5),
     RGBColor(0xE8, 0xF0, 0xFE))
tb(s, Inches(0.6), Inches(1.65), Inches(5.8), Inches(0.5),
   '✅ 이미 완료된 부분', size=16, bold=True, color=C_ACCENT)
tb(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.7),
   [
       '☑ 엑셀 v3 양식 (33컬럼, 4시트)',
       '☑ 데이터 유효성 검사 (드롭다운)',
       '☑ 색상범례 / 사용가이드 시트',
       '☑ 등록 페이지 (/admin/register)',
       '☑ 메타 4영역 (출처/코드/분류/속성)',
       '☑ 블록 시퀀스 입력 UI',
       '☑ 실시간 HTML 미리보기',
       '☑ JSON 출력 / 다운로드',
       '☑ 양방향 변환 (xlsx ↔ json)',
       '☑ 왕복 일치 검증 ✅',
   ], size=13)

rect(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.5),
     RGBColor(0xE8, 0xF8, 0xE8))
tb(s, Inches(7.0), Inches(1.65), Inches(5.8), Inches(0.5),
   '📋 다음 단계 (개발자 작업)', size=16, bold=True, color=C_OK)
tb(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(4.7),
   [
       '☐ 저장 백엔드 결정 (DB / 파일 시스템)',
       '☐ 등록 페이지 → DB 저장 API',
       '☐ HTML sanitize (DOMPurify / bleach)',
       '   - 허용 화이트리스트 적용',
       '   - <script> <iframe> on* 제거',
       '☐ 그림 파일 업로드 (Storage)',
       '☐ [필요: …] 마커 검수 큐',
       '☐ 시험 페이지 (/cbt) 데이터 모델 연결',
       '   - 문항코드 순 정렬',
       '   - 회차/연도/빈출도 필터',
       '☐ 일괄 등록 (엑셀 업로드 → DB import)',
       '☐ 모바일 렌더링 검증',
       '   - <table> 좁은 화면 처리',
   ], size=13)

prs.save(OUT)
print(f'WROTE: {OUT}')
print(f'  슬라이드 {len(prs.slides)}장')
