"""사용자 수정 PPT 슬라이드 내용 추출 — 텍스트/제목/순서 확인."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

PATH = 'data/presentation/엑셀-등록페이지_연동.pptx'
prs = Presentation(PATH)
print(f'슬라이드 {len(prs.slides)}장')

for i, slide in enumerate(prs.slides, 1):
    print(f'\n===== Slide {i} =====')
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                lines = [l for l in txt.splitlines() if l.strip()]
                head = lines[0] if lines else ''
                if len(head) > 80:
                    head = head[:80] + '…'
                print(f'  · {head}')
                if len(lines) > 1:
                    for line in lines[1:6]:
                        if len(line) > 70:
                            line = line[:70] + '…'
                        print(f'      {line}')
        elif shape.shape_type == 13:  # picture
            print(f'  [그림]')
